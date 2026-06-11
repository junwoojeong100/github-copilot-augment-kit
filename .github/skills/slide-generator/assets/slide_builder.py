"""Build readable, design-rich PowerPoint decks from a structured JSON slide spec.

Supports a Microsoft-inspired card-based visual system: eyebrow labels,
page badges, multi-color accent series, KPI stats, timelines, card grids,
and optional insight/source footnote boxes.
"""

from __future__ import annotations

import argparse
import json
from collections.abc import Iterable, Mapping, Sequence
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

THEME: dict[str, Any] = {
    "font": "Segoe UI",          # latin font
    "kr_font": "Malgun Gothic",  # east-asian (Korean) font
    "title_font_size": 40,
    "subtitle_font_size": 20,
    "heading_font_size": 26,
    "eyebrow_font_size": 12,
    "body_font_size": 17,
    "small_font_size": 10,
    "colors": {
        "primary": "0F2A4A",
        "primary_soft": "16395F",
        "accent": "0078D4",
        "accent_dark": "005A9E",
        "gold": "E6A532",
        "teal": "0E7C86",
        "purple": "6B4FA1",
        "background": "FFFFFF",
        "section_background": "F3F7FC",
        "card": "FFFFFF",
        "card_alt": "F3F7FC",
        "border": "D8E0EA",
        "mid_gray": "9AB4CF",
        "text": "1F2937",
        "muted_text": "5B677A",
        "white": "FFFFFF",
    },
    # accent series used to color-cycle cards / markers / kpi
    "series": ["accent", "gold", "teal", "purple"],
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_X = Inches(0.72)
CONTENT_W = Inches(11.9)
HEAD_TOP = Inches(0.52)
CONTENT_BOTTOM = Inches(6.78)
FOOTER_TOP = Inches(7.05)
MAX_BULLETS = 6
MAX_NESTED_BULLETS = 3

SlideSpec = Mapping[str, Any]


class SlideSpecError(ValueError):
    """Raised when the deck specification is invalid."""


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def build_presentation(spec: Mapping[str, Any]) -> Presentation:
    """Create a PowerPoint presentation from a validated JSON-like spec."""
    theme = _merged_theme(spec.get("theme"))
    _validate_spec(spec)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    deck_title = str(spec["title"])
    slides = spec["slides"]
    total = len(slides)
    dispatch = {
        "title": _add_title_slide,
        "section": _add_section_slide,
        "bullets": _add_bullets_slide,
        "two_column": _add_two_column_slide,
        "table": _add_table_slide,
        "callout": _add_callout_slide,
        "closing": _add_closing_slide,
        "cards": _add_cards_slide,
        "kpi": _add_kpi_slide,
        "timeline": _add_timeline_slide,
        "architecture": _add_architecture_slide,
        "process": _add_process_slide,
    }
    for index, slide_spec in enumerate(slides, start=1):
        slide_type = slide_spec["type"]
        builder = dispatch.get(slide_type)
        if builder is None:
            raise SlideSpecError(f"Unsupported slide type: {slide_type}")
        if slide_type == "title":
            builder(prs, slide_spec, spec, theme, index, total)
        else:
            slide = builder(prs, slide_spec, theme, index, total)
            dark_callout = (slide_type == "callout"
                            and str(slide_spec.get("style", "")).lower() == "dark")
            if slide_type not in {"closing"} and not dark_callout:
                _add_footer(slide, deck_title, index, total, theme)
    return prs


def build_pptx(spec_path: str | Path, output_path: str | Path) -> Path:
    """Read a JSON spec and write a .pptx file."""
    spec_file = Path(spec_path)
    output_file = Path(output_path)
    with spec_file.open("r", encoding="utf-8") as handle:
        spec = json.load(handle)

    prs = build_presentation(spec)
    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_file)
    return output_file


# ---------------------------------------------------------------------------
# Validation
# ---------------------------------------------------------------------------

def _validate_spec(spec: Mapping[str, Any]) -> None:
    if not isinstance(spec, Mapping):
        raise SlideSpecError("Spec must be a JSON object.")
    if not spec.get("title"):
        raise SlideSpecError("Spec requires a non-empty 'title'.")
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        raise SlideSpecError("Spec requires a non-empty 'slides' array.")

    known = {"title", "section", "bullets", "two_column", "table",
             "callout", "closing", "cards", "kpi", "timeline",
             "architecture", "process"}
    for index, slide in enumerate(slides, start=1):
        if not isinstance(slide, Mapping):
            raise SlideSpecError(f"Slide {index} must be an object.")
        slide_type = slide.get("type")
        if not slide_type:
            raise SlideSpecError(f"Slide {index} requires a 'type'.")
        if slide_type not in known:
            raise SlideSpecError(f"Slide {index} has unknown type '{slide_type}'.")

        needs_heading = {"bullets", "two_column", "table", "callout",
                         "closing", "cards", "kpi", "timeline",
                         "architecture", "process"}
        if slide_type in needs_heading and not slide.get("heading"):
            raise SlideSpecError(f"Slide {index} ({slide_type}) requires 'heading'.")
        if slide_type == "bullets" and not isinstance(slide.get("bullets"), list):
            raise SlideSpecError(f"Slide {index} (bullets) requires a 'bullets' array.")
        if slide_type == "two_column":
            for side in ("left", "right"):
                column = slide.get(side)
                if not isinstance(column, Mapping) or not isinstance(column.get("bullets"), list):
                    raise SlideSpecError(f"Slide {index} (two_column) requires '{side}.bullets'.")
        if slide_type == "cards":
            cards = slide.get("cards")
            if not isinstance(cards, list) or not cards:
                raise SlideSpecError(f"Slide {index} (cards) requires a non-empty 'cards' array.")
        if slide_type == "kpi":
            items = slide.get("items")
            if not isinstance(items, list) or not items:
                raise SlideSpecError(f"Slide {index} (kpi) requires a non-empty 'items' array.")
        if slide_type == "timeline":
            phases = slide.get("phases")
            if not isinstance(phases, list) or not phases:
                raise SlideSpecError(f"Slide {index} (timeline) requires a non-empty 'phases' array.")
        if slide_type == "table":
            columns = slide.get("columns")
            rows = slide.get("rows")
            if not isinstance(columns, list) or not columns:
                raise SlideSpecError(f"Slide {index} (table) requires non-empty 'columns'.")
            if not isinstance(rows, list) or not rows:
                raise SlideSpecError(f"Slide {index} (table) requires non-empty 'rows'.")
            for row_number, row in enumerate(rows, start=1):
                if not isinstance(row, list):
                    raise SlideSpecError(f"Slide {index} table row {row_number} must be an array.")
        if slide_type == "callout" and not slide.get("text"):
            raise SlideSpecError(f"Slide {index} (callout) requires 'text'.")
        if slide_type == "architecture":
            layers = slide.get("layers")
            if not isinstance(layers, list) or not layers:
                raise SlideSpecError(f"Slide {index} (architecture) requires a non-empty 'layers' array.")
        if slide_type == "process":
            steps = slide.get("steps")
            if not isinstance(steps, list) or not steps:
                raise SlideSpecError(f"Slide {index} (process) requires a non-empty 'steps' array.")


def _merged_theme(custom_theme: Any) -> dict[str, Any]:
    theme = {**THEME, "colors": dict(THEME["colors"]), "series": list(THEME["series"])}
    if isinstance(custom_theme, Mapping):
        for key, value in custom_theme.items():
            if key == "colors" and isinstance(value, Mapping):
                theme["colors"].update({str(k): str(v).lstrip("#") for k, v in value.items()})
            elif key == "series" and isinstance(value, list):
                theme["series"] = [str(v) for v in value]
            else:
                theme[key] = value
    return theme


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _add_title_slide(prs, slide_spec, deck_spec, theme, index, total) -> None:
    slide = _blank_slide(prs, theme)
    _set_background(slide, "primary", theme)

    # decorative side band + accent column
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.26), SLIDE_HEIGHT)
    _fill(band, theme, "gold")
    _line_none(band)
    _corner_accents(slide, theme)

    eyebrow = str(slide_spec.get("eyebrow") or deck_spec.get("eyebrow") or "")
    title = str(slide_spec.get("heading") or deck_spec.get("title", ""))
    subtitle = str(slide_spec.get("subtitle") or deck_spec.get("subtitle") or "")
    author = str(slide_spec.get("author") or deck_spec.get("author") or "")

    if eyebrow:
        _add_text(slide, eyebrow.upper(), Inches(0.95), Inches(1.35), Inches(10.8), Inches(0.4),
                  theme, "gold", int(theme["eyebrow_font_size"]) + 1, bold=True, spacing=2.0)
    _add_text(slide, title, Inches(0.92), Inches(1.95), Inches(11.2), Inches(1.7),
              theme, "white", int(theme["title_font_size"]), bold=True)
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.96), Inches(3.5), Inches(2.2), Inches(0.05))
    _fill(divider, theme, "gold")
    _line_none(divider)
    if subtitle:
        _add_text(slide, subtitle, Inches(0.96), Inches(3.7), Inches(10.8), Inches(0.85),
                  theme, "mid_gray", int(theme["subtitle_font_size"]))

    chips = slide_spec.get("chips") or deck_spec.get("chips")
    if isinstance(chips, list) and chips:
        _add_chip_row(slide, [str(c) for c in chips], Inches(0.96), Inches(5.05),
                      theme, on_dark=True)
    if author:
        _add_text(slide, author, Inches(0.96), Inches(6.55), Inches(9.0), Inches(0.4),
                  theme, "mid_gray", 13)


def _add_section_slide(prs, slide_spec, theme, index, total):
    slide = _blank_slide(prs, theme)
    _set_background(slide, "section_background", theme)
    accent = _series_color(theme, index)

    num = slide_spec.get("number")
    if num is not None:
        _add_text(slide, str(num), Inches(0.95), Inches(1.6), Inches(3.0), Inches(1.6),
                  theme, accent, 96, bold=True)
    eyebrow = slide_spec.get("eyebrow")
    if eyebrow:
        _add_text(slide, str(eyebrow).upper(), Inches(1.0), Inches(3.05), Inches(10.9), Inches(0.4),
                  theme, accent, int(theme["eyebrow_font_size"]), bold=True, spacing=2.0)
    _add_text(slide, str(slide_spec.get("heading", "Section")), Inches(0.98), Inches(3.5),
              Inches(11.2), Inches(1.2), theme, "primary", 34, bold=True)
    subtitle = slide_spec.get("subtitle")
    if subtitle:
        _add_text(slide, str(subtitle), Inches(1.02), Inches(4.75), Inches(10.4), Inches(0.6),
                  theme, "muted_text", 18)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.5), Inches(2.4), Inches(0.1))
    _fill(bar, theme, accent)
    _line_none(bar)
    return slide


def _add_bullets_slide(prs, slide_spec, theme, index, total):
    slide, top, height = _content_slide(prs, slide_spec, theme, index, total)
    bullets = list(slide_spec.get("bullets", []))[:MAX_BULLETS]
    if not bullets:
        bullets = ["내용을 추가하세요."]
    n = len(bullets)
    gap = Inches(0.16)
    row_h = (height - gap * (n - 1)) / n
    y = top
    for i, bullet in enumerate(bullets):
        color = _series_color(theme, i)
        _add_bullet_card(slide, bullet, MARGIN_X, y, CONTENT_W, row_h, theme, color)
        y = y + row_h + gap
    return slide


def _add_two_column_slide(prs, slide_spec, theme, index, total):
    slide, top, height = _content_slide(prs, slide_spec, theme, index, total)
    left = slide_spec.get("left", {})
    right = slide_spec.get("right", {})
    gap = Inches(0.4)
    width = (CONTENT_W - gap) / 2
    _add_column_card(slide, left, MARGIN_X, top, width, height, theme,
                     left.get("color") or _series_color(theme, 0))
    _add_column_card(slide, right, MARGIN_X + width + gap, top, width, height, theme,
                     right.get("color") or _series_color(theme, 2))
    return slide


def _add_cards_slide(prs, slide_spec, theme, index, total):
    slide, top, height = _content_slide(prs, slide_spec, theme, index, total)
    cards = list(slide_spec.get("cards", []))[:6]
    n = len(cards)
    cols = slide_spec.get("columns")
    if not isinstance(cols, int) or cols < 1:
        cols = 3 if n in (3, 5, 6) else (2 if n in (2, 4) else min(n, 3))
    cols = min(cols, n)
    rows = (n + cols - 1) // cols
    numbered = bool(slide_spec.get("numbered", False))
    gap = Inches(0.35)
    cw = (CONTENT_W - gap * (cols - 1)) / cols
    ch = (height - gap * (rows - 1)) / rows
    for i, card in enumerate(cards):
        r, c = divmod(i, cols)
        x = MARGIN_X + c * (cw + gap)
        y = top + r * (ch + gap)
        _add_feature_card(slide, card, x, y, cw, ch, theme, _series_color(theme, i),
                          i + 1, numbered=numbered)
    return slide


def _add_kpi_slide(prs, slide_spec, theme, index, total):
    slide, top, height = _content_slide(prs, slide_spec, theme, index, total)
    items = list(slide_spec.get("items", []))[:4]
    n = len(items)
    solid = str(slide_spec.get("style", "")).lower() == "solid"
    gap = Inches(0.35)
    cw = (CONTENT_W - gap * (n - 1)) / n
    card_h = height if height < Inches(2.7) else Inches(2.7)
    for i, item in enumerate(items):
        x = MARGIN_X + i * (cw + gap)
        _add_kpi_card(slide, item, x, top, cw, card_h, theme, _series_color(theme, i), solid=solid)
    return slide


def _add_timeline_slide(prs, slide_spec, theme, index, total):
    slide, top, height = _content_slide(prs, slide_spec, theme, index, total)
    phases = list(slide_spec.get("phases", []))[:4]
    n = len(phases)
    gap = Inches(0.35)
    cw = (CONTENT_W - gap * (n - 1)) / n
    # connecting line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X + Inches(0.2), top + Inches(0.34),
                                  CONTENT_W - Inches(0.4), Inches(0.04))
    _fill(line, theme, "border")
    _line_none(line)
    for i, phase in enumerate(phases):
        x = MARGIN_X + i * (cw + gap)
        _add_timeline_card(slide, phase, x, top, cw, height, theme, _series_color(theme, i), i + 1)
    return slide


def _add_architecture_slide(prs, slide_spec, theme, index, total):
    """Layered architecture diagram: optional left sidebar + stacked layers
    (one may be highlighted) + optional full-width footer governance bar."""
    slide, top, height = _content_slide(prs, slide_spec, theme, index, total)
    sidebar = slide_spec.get("sidebar")
    layers = list(slide_spec.get("layers", []))[:5]
    footer = slide_spec.get("footer")

    footer_h = Inches(0.78) if isinstance(footer, Mapping) else Inches(0.0)
    footer_gap = Inches(0.2) if isinstance(footer, Mapping) else Inches(0.0)
    stack_h = height - footer_h - footer_gap

    if isinstance(sidebar, Mapping):
        sb_w = Inches(2.7)
        _panel(slide, MARGIN_X, top, sb_w, stack_h, theme, "primary",
               line_color=None, shadow=True, radius=0.11)
        sx = MARGIN_X + Inches(0.28)
        sw = sb_w - Inches(0.56)
        _add_text(slide, str(sidebar.get("eyebrow") or sidebar.get("title", "")),
                  sx, top + Inches(0.3), sw, Inches(0.4), theme, "gold", 14, bold=True)
        _add_text(slide, str(sidebar.get("name") or sidebar.get("subtitle", "")),
                  sx, top + Inches(0.78), sw, Inches(0.6), theme, "white", 16, bold=True)
        if sidebar.get("text"):
            _add_text(slide, str(sidebar["text"]), sx, top + Inches(1.5), sw, stack_h - Inches(1.7),
                      theme, "mid_gray", 12)
        layers_x = MARGIN_X + sb_w + Inches(0.3)
        layers_w = CONTENT_W - sb_w - Inches(0.3)
    else:
        layers_x = MARGIN_X
        layers_w = CONTENT_W

    n = len(layers)
    gap = Inches(0.18)
    lh = (stack_h - gap * (n - 1)) / n
    ly = top
    for i, layer in enumerate(layers):
        highlight = isinstance(layer, Mapping) and layer.get("highlight")
        fill_name = "accent" if highlight else "card_alt"
        title_color = "white" if highlight else "primary"
        desc_color = "white" if highlight else "muted_text"
        _panel(slide, layers_x, ly, layers_w, lh, theme, fill_name,
               line_color=None if highlight else "border", shadow=True, radius=0.09)
        if isinstance(layer, Mapping):
            title = str(layer.get("title", ""))
            desc = str(layer.get("text") or layer.get("desc") or "")
        else:
            title, desc = str(layer), ""
        pad = Inches(0.32)
        text_w = layers_w - Inches(0.64)
        if not desc:
            _add_text(slide, title, layers_x + pad, ly, text_w, lh, theme, title_color,
                      15, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        elif lh >= Inches(0.82):
            # tall layer: stack title over description
            _add_text(slide, title, layers_x + pad, ly + Inches(0.12), text_w, Inches(0.4),
                      theme, title_color, 15, bold=True)
            _add_text(slide, desc, layers_x + pad, ly + Inches(0.52), text_w, lh - Inches(0.58),
                      theme, desc_color, 12)
        else:
            # short layer: title left, description right, both vertically centered
            title_w = int(layers_w * 0.42)
            _add_text(slide, title, layers_x + pad, ly, title_w, lh, theme, title_color,
                      15, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            _add_text(slide, desc, layers_x + pad + title_w, ly, text_w - title_w, lh,
                      theme, desc_color, 12, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                               layers_x + layers_w / 2 - Inches(0.015),
                                               ly + lh, Inches(0.03), gap)
            _fill(connector, theme, "mid_gray")
            _line_none(connector)
        ly = ly + lh + gap

    if isinstance(footer, Mapping):
        fy = top + stack_h + footer_gap
        _panel(slide, MARGIN_X, fy, CONTENT_W, footer_h, theme, "primary_soft",
               line_color=None, shadow=True, radius=0.09)
        _add_text(slide, str(footer.get("title", "")), MARGIN_X + Inches(0.32), fy,
                  Inches(3.4), footer_h, theme, "gold", 13, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        _add_text(slide, str(footer.get("text", "")), MARGIN_X + Inches(3.7), fy,
                  CONTENT_W - Inches(4.0), footer_h, theme, "white", 12, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def _add_process_slide(prs, slide_spec, theme, index, total):
    """Horizontal step flow with chevron connectors between boxes."""
    slide, top, height = _content_slide(prs, slide_spec, theme, index, total)
    steps = list(slide_spec.get("steps", []))[:6]
    n = len(steps)
    arrow_w = Inches(0.5)
    bw = (CONTENT_W - arrow_w * (n - 1)) / n
    has_text = any(isinstance(s, Mapping) and s.get("text") for s in steps)
    box_h = Inches(1.7) if has_text else Inches(1.1)
    by = top + (height - box_h) / 2
    x = MARGIN_X
    for i, step in enumerate(steps):
        color = "primary" if i % 2 == 0 else "accent_dark"
        _panel(slide, x, by, bw, box_h, theme, color, line_color=None, shadow=True, radius=0.1)
        if isinstance(step, Mapping):
            label = str(step.get("label") or step.get("title", ""))
            text = str(step.get("text", ""))
        else:
            label, text = str(step), ""
        if text:
            _add_text(slide, label, x + Inches(0.2), by + Inches(0.22), bw - Inches(0.4), Inches(0.5),
                      theme, "white", 15, bold=True, align=PP_ALIGN.CENTER)
            _add_text(slide, text, x + Inches(0.2), by + Inches(0.78), bw - Inches(0.4),
                      box_h - Inches(0.95), theme, "mid_gray", 11, align=PP_ALIGN.CENTER)
        else:
            _add_text(slide, label, x + Inches(0.15), by, bw - Inches(0.3), box_h, theme, "white",
                      15, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x = x + bw
        if i < n - 1:
            _add_text(slide, "›", x, by, arrow_w, box_h, theme, "gold", 26, bold=True,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x = x + arrow_w
    return slide


def _add_table_slide(prs, slide_spec, theme, index, total):
    slide, top, height = _content_slide(prs, slide_spec, theme, index, total)
    columns = [str(column) for column in slide_spec.get("columns", [])]
    rows = _normalized_rows(slide_spec.get("rows", []), len(columns))
    max_rows = rows[:7]
    if len(rows) > len(max_rows):
        max_rows[-1] = [*max_rows[-1][:-1], f"{max_rows[-1][-1]} …"]
    emphasize_first = bool(slide_spec.get("emphasize_first_column", True))

    table_shape = slide.shapes.add_table(len(max_rows) + 1, len(columns), MARGIN_X, top, CONTENT_W, height)
    table = table_shape.table
    _strip_table_style(table)
    table.first_row = False
    for col_idx in range(len(columns)):
        table.columns[col_idx].width = int(CONTENT_W / len(columns))

    for col_idx, column in enumerate(columns):
        cell = table.cell(0, col_idx)
        _style_cell(cell, column, theme, "primary", "white", bold=True, size=12,
                    align=PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER)

    for row_idx, row in enumerate(max_rows, start=1):
        fill_color = "card_alt" if row_idx % 2 == 0 else "white"
        for col_idx, value in enumerate(row):
            first = emphasize_first and col_idx == 0
            _style_cell(table.cell(row_idx, col_idx), value, theme, fill_color,
                        "primary" if first else "text",
                        bold=first, size=12 if first else 11,
                        align=PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER)
    return slide


def _add_callout_slide(prs, slide_spec, theme, index, total):
    if str(slide_spec.get("style", "")).lower() == "dark":
        return _add_callout_dark(prs, slide_spec, theme)
    slide, top, height = _content_slide(prs, slide_spec, theme, index, total)
    accent = _series_color(theme, index)
    _panel(slide, MARGIN_X, top + Inches(0.2), CONTENT_W, Inches(2.7), theme,
           "section_background", shadow=True, radius=0.12)
    _bar(slide, MARGIN_X + Inches(0.34), top + Inches(0.6), Inches(0.1), Inches(1.9), theme, accent)
    _add_text(slide, str(slide_spec["text"]), MARGIN_X + Inches(0.7), top + Inches(0.5),
              CONTENT_W - Inches(1.2), Inches(2.1), theme, "primary", 26, bold=True,
              anchor=MSO_ANCHOR.MIDDLE)
    chips = slide_spec.get("chips")
    if isinstance(chips, list) and chips:
        _add_chip_row(slide, [str(c) for c in chips], MARGIN_X, top + Inches(3.25), theme)
    return slide


def _add_callout_dark(prs, slide_spec, theme):
    """Full navy 'big statement' slide (e.g. KEY MESSAGE), no content header."""
    slide = _blank_slide(prs, theme)
    _set_background(slide, "primary", theme)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.26), SLIDE_HEIGHT)
    _fill(band, theme, "gold")
    _line_none(band)

    eyebrow = str(slide_spec.get("eyebrow") or "KEY MESSAGE")
    _add_text(slide, eyebrow.upper(), Inches(1.1), Inches(1.15), Inches(11.0), Inches(0.4),
              theme, "gold", 14, bold=True, spacing=2.0)
    seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.12), Inches(1.75), Inches(1.0), Inches(0.06))
    _fill(seg, theme, "gold")
    _line_none(seg)
    _add_text(slide, str(slide_spec["text"]), Inches(1.1), Inches(2.1), Inches(11.1), Inches(2.9),
              theme, "white", 27, bold=True)
    note = slide_spec.get("note")
    if note:
        _add_text(slide, str(note), Inches(1.12), Inches(5.0), Inches(11.0), Inches(1.0),
                  theme, "mid_gray", 15)
    chips = slide_spec.get("chips")
    if isinstance(chips, list) and chips:
        _add_chip_row(slide, [str(c) for c in chips], Inches(1.1), Inches(6.2), theme, on_dark=True)
    return slide


def _add_closing_slide(prs, slide_spec, theme, index, total):
    slide = _blank_slide(prs, theme)
    _set_background(slide, "primary", theme)
    _corner_accents(slide, theme)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.26), SLIDE_HEIGHT)
    _fill(band, theme, "gold")
    _line_none(band)

    _add_text(slide, str(slide_spec.get("heading", "Thank you")), Inches(0.95), Inches(2.5),
              Inches(11.4), Inches(1.2), theme, "white", 38, bold=True, align=PP_ALIGN.CENTER)
    contact = slide_spec.get("contact")
    if contact:
        _add_text(slide, str(contact), Inches(1.4), Inches(3.95), Inches(10.5), Inches(0.9),
                  theme, "mid_gray", 16, align=PP_ALIGN.CENTER)
    chips = slide_spec.get("chips")
    if isinstance(chips, list) and chips:
        _add_chip_row(slide, [str(c) for c in chips], Inches(0), Inches(5.2),
                      theme, on_dark=True, center=True)
    return slide


# ---------------------------------------------------------------------------
# Composite components
# ---------------------------------------------------------------------------

def _content_slide(prs, slide_spec, theme, index, total):
    """Render the shared content header (eyebrow + heading + underline accent +
    page badge) and optional insight/source footnote boxes.
    Returns (slide, content_top, content_height)."""
    slide = _blank_slide(prs, theme)
    accent = _series_color(theme, index)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.14), SLIDE_HEIGHT)
    _fill(bar, theme, accent)
    _line_none(bar)

    eyebrow = slide_spec.get("eyebrow")
    head_y = HEAD_TOP
    if eyebrow:
        _add_text(slide, str(eyebrow).upper(), MARGIN_X, head_y, CONTENT_W, Inches(0.32),
                  theme, accent, int(theme["eyebrow_font_size"]), bold=True, spacing=1.8)
        head_y = head_y + Inches(0.36)
    _add_text(slide, str(slide_spec["heading"]), MARGIN_X, head_y, Inches(10.2), Inches(0.66),
              theme, "primary", int(theme["heading_font_size"]), bold=True)

    # page badge top-right
    _page_badge(slide, index, total, theme, accent)

    # underline accent: full-width hairline + short accent segment over it
    rule_y = head_y + Inches(0.64)
    hair = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X, rule_y, CONTENT_W, Inches(0.015))
    _fill(hair, theme, "border")
    _line_none(hair)
    seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X, rule_y - Inches(0.006),
                                 Inches(1.5), Inches(0.034))
    _fill(seg, theme, accent)
    _line_none(seg)

    top = rule_y + Inches(0.24)
    subhead = slide_spec.get("subhead")
    if subhead:
        _add_text(slide, str(subhead), MARGIN_X, top, CONTENT_W, Inches(0.34),
                  theme, "muted_text", 13)
        top = top + Inches(0.46)

    bottom = CONTENT_BOTTOM
    insight = slide_spec.get("insight")
    source = slide_spec.get("source")
    if insight or source:
        box_top = Inches(6.05)
        bottom = box_top - Inches(0.15)
        _footnote_boxes(slide, insight, source, box_top, theme, accent)
    return slide, top, bottom - top


def _page_badge(slide, index, total, theme, accent):
    w, h = Inches(1.0), Inches(0.42)
    x = SLIDE_WIDTH - MARGIN_X - w
    _panel(slide, x, HEAD_TOP + Inches(0.02), w, h, theme, "section_background",
           shadow=False, radius=0.07)
    _add_text(slide, f"{index:02d} / {total:02d}", x, HEAD_TOP + Inches(0.04), w, Inches(0.38),
              theme, "muted_text", 11, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _footnote_boxes(slide, insight, source, y, theme, accent):
    if insight and source:
        w1 = Inches(7.7)
        w2 = CONTENT_W - w1 - Inches(0.3)
        _kv_box(slide, "핵심", str(insight), MARGIN_X, y, w1, Inches(0.62), theme, accent)
        _kv_box(slide, "출처", str(source), MARGIN_X + w1 + Inches(0.3), y, w2, Inches(0.62),
                theme, "mid_gray", muted=True)
    elif insight:
        _kv_box(slide, "핵심", str(insight), MARGIN_X, y, CONTENT_W, Inches(0.62), theme, accent)
    elif source:
        _kv_box(slide, "출처", str(source), MARGIN_X, y, CONTENT_W, Inches(0.62), theme, "mid_gray", muted=True)


def _kv_box(slide, label, value, x, y, w, h, theme, accent, muted=False):
    _panel(slide, x, y, w, h, theme, "card_alt" if not muted else "background",
           shadow=True, radius=0.07)
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.12), y + Inches(0.13),
                                 Inches(0.62), h - Inches(0.26))
    _fill(tag, theme, accent)
    _line_none(tag)
    _rounded(tag, 0.05)
    _add_text(slide, label, x + Inches(0.12), y + Inches(0.13), Inches(0.62), h - Inches(0.26),
              theme, "white", 11, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(slide, value, x + Inches(0.86), y, w - Inches(1.0), h, theme,
              "text" if not muted else "muted_text", 11, anchor=MSO_ANCHOR.MIDDLE)


def _add_bullet_card(slide, bullet, x, y, w, h, theme, accent):
    _panel(slide, x, y, w, h, theme, "card_alt", shadow=True, radius=0.1)
    _bar(slide, x + Inches(0.18), y + Inches(0.16), Inches(0.09), h - Inches(0.32), theme, accent)

    if isinstance(bullet, Mapping):
        text = str(bullet.get("text") or bullet.get("title") or "")
        children = list(bullet.get("children") or bullet.get("bullets") or [])[:MAX_NESTED_BULLETS]
    else:
        text, children = str(bullet), []

    tx = x + Inches(0.45)
    tw = w - Inches(0.7)
    if children:
        _add_text(slide, text, tx, y + Inches(0.12), tw, Inches(0.42), theme, "primary",
                  int(theme["body_font_size"]), bold=True)
        sub = "   •   ".join(str(c) for c in children)
        _add_text(slide, sub, tx, y + Inches(0.55), tw, h - Inches(0.65), theme, "muted_text", 13)
    else:
        _add_text(slide, text, tx, y, tw, h, theme, "text", int(theme["body_font_size"]),
                  anchor=MSO_ANCHOR.MIDDLE)


def _add_column_card(slide, column, x, y, width, height, theme, accent):
    _panel(slide, x, y, width, height, theme, "card", shadow=True, radius=0.1)
    _panel(slide, x, y, width, Inches(0.72), theme, accent,
           line_color=None, shadow=False, radius=0.1)
    _add_text(slide, str(column.get("title", "")), x + Inches(0.3), y, width - Inches(0.6),
              Inches(0.72), theme, "white", 16, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _add_bullets(slide, column.get("bullets", []), x + Inches(0.3), y + Inches(0.92),
                 width - Inches(0.6), height - Inches(1.1), theme, accent, font_size=15)


def _add_feature_card(slide, card, x, y, w, h, theme, accent, number, numbered=False):
    _panel(slide, x, y, w, h, theme, "card", shadow=True, radius=0.1)
    _bar(slide, x, y + Inches(0.16), Inches(0.1), h - Inches(0.32), theme, accent)

    if isinstance(card, Mapping):
        title = str(card.get("title", ""))
        desc = str(card.get("text") or card.get("desc") or "")
        badge = str(card.get("badge") or number)
    else:
        title, desc, badge = str(card), "", str(number)

    inner_x = x + Inches(0.42)
    inner_w = w - Inches(0.66)
    if numbered:
        bsize = Inches(0.5)
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inner_x, y + Inches(0.3),
                                      bsize, bsize)
        _fill(chip, theme, accent)
        _line_none(chip)
        _add_text(slide, badge, inner_x, y + Inches(0.3), bsize, bsize, theme, "white",
                  16, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text(slide, title, inner_x + Inches(0.66), y + Inches(0.3), inner_w - Inches(0.66),
                  Inches(0.5), theme, "primary", 15, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        div_y = y + Inches(0.96)
    else:
        _add_text(slide, title, inner_x, y + Inches(0.32), inner_w, Inches(0.52),
                  theme, "primary", 15, bold=True)
        div_y = y + Inches(0.92)

    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inner_x, div_y, inner_w, Inches(0.014))
    _fill(divider, theme, "border")
    _line_none(divider)
    if desc:
        _add_text(slide, desc, inner_x, div_y + Inches(0.14), inner_w,
                  h - (div_y - y) - Inches(0.3), theme, "muted_text", 13)


def _add_kpi_card(slide, item, x, y, w, h, theme, accent, solid=False):
    if solid:
        _panel(slide, x, y, w, h, theme, "primary", line_color=None, shadow=True, radius=0.1)
        value_color, label_color, cap_color = "gold", "white", "mid_gray"
    else:
        _panel(slide, x, y, w, h, theme, "card", shadow=True, radius=0.1)
        _bar(slide, x, y + Inches(0.18), Inches(0.1), h - Inches(0.36), theme, accent)
        value_color, label_color, cap_color = accent, "primary", "muted_text"

    value = str(item.get("value", "")) if isinstance(item, Mapping) else str(item)
    label = str(item.get("label", "")) if isinstance(item, Mapping) else ""
    caption = str(item.get("caption", "")) if isinstance(item, Mapping) else ""

    value_size = 40 if len(value) <= 6 else (30 if len(value) <= 8 else 23)
    _add_text(slide, value, x + Inches(0.34), y + Inches(0.42), w - Inches(0.5), Inches(1.15),
              theme, value_color, value_size, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(slide, label, x + Inches(0.36), y + Inches(1.62), w - Inches(0.6), Inches(0.5),
              theme, label_color, 15, bold=True)
    if caption:
        _add_text(slide, caption, x + Inches(0.36), y + Inches(2.08), w - Inches(0.6), Inches(0.55),
                  theme, cap_color, 11)


def _add_timeline_card(slide, phase, x, y, w, h, theme, accent, number):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.05), y + Inches(0.12),
                                 Inches(0.48), Inches(0.48))
    _fill(dot, theme, accent)
    _line_none(dot)
    _soft_shadow(dot)
    _add_text(slide, str(number), x + Inches(0.05), y + Inches(0.12), Inches(0.48), Inches(0.48),
              theme, "white", 16, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _panel(slide, x, y + Inches(0.85), w, h - Inches(0.85), theme, "card",
           shadow=True, radius=0.1)

    label = str(phase.get("label", "")) if isinstance(phase, Mapping) else str(phase)
    title = str(phase.get("title", "")) if isinstance(phase, Mapping) else ""
    items = list(phase.get("items", [])) if isinstance(phase, Mapping) else []

    _add_text(slide, label, x + Inches(0.28), y + Inches(1.0), w - Inches(0.5), Inches(0.4),
              theme, accent, 15, bold=True)
    if title:
        _add_text(slide, title, x + Inches(0.28), y + Inches(1.42), w - Inches(0.5), Inches(0.55),
                  theme, "primary", 13, bold=True)
    if items:
        _add_bullets(slide, items, x + Inches(0.28), y + Inches(2.0), w - Inches(0.5),
                     h - Inches(2.2), theme, accent, font_size=12)


def _add_chip_row(slide, chips, x, y, theme, on_dark=False, center=False):
    chips = chips[:6]
    pad = Inches(0.22)
    cur = MARGIN_X if center else x
    for chip in chips:
        cw = Inches(0.5) + Inches(0.105) * min(len(chip), 22)
        if on_dark:
            _panel(slide, cur, y, cw, Inches(0.46), theme, "primary_soft",
                   line_color="accent", shadow=False, radius=0.08)
            tcolor = "white"
        else:
            _panel(slide, cur, y, cw, Inches(0.46), theme, "section_background",
                   line_color="border", shadow=False, radius=0.08)
            tcolor = "primary"
        _add_text(slide, chip, cur, y, cw, Inches(0.46), theme, tcolor, 12, bold=True,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cur = cur + cw + pad


def _add_bullets(slide, bullets, x, y, width, height, theme, accent=None, font_size=None):
    box = slide.shapes.add_textbox(x, y, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.05)
    frame.word_wrap = True
    size = font_size or int(theme["body_font_size"])

    rendered = list(_flatten_bullets(bullets))[:MAX_BULLETS + MAX_BULLETS * MAX_NESTED_BULLETS]
    if not rendered:
        rendered = [(0, "내용을 추가하세요.")]

    for idx, (level, text) in enumerate(rendered):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        prefix = "•" if level == 0 else "–"
        paragraph.text = f"{prefix} {text}"
        paragraph.level = min(level, 2)
        paragraph.space_after = Pt(6 if level == 0 else 3)
        _style_paragraph(paragraph, theme,
                         size if level == 0 else max(size - 2, 11),
                         "text" if level == 0 else "muted_text")


def _flatten_bullets(bullets):
    top_count = 0
    for bullet in bullets:
        if top_count >= MAX_BULLETS:
            yield 0, "핵심 항목이 더 있습니다. 상세 내용은 appendix 또는 발표자 노트로 분리하세요."
            return
        top_count += 1
        if isinstance(bullet, Mapping):
            text = str(bullet.get("text") or bullet.get("title") or "")
            if text:
                yield 0, text
            for sub in list(bullet.get("children") or bullet.get("bullets") or [])[:MAX_NESTED_BULLETS]:
                yield 1, str(sub)
        else:
            yield 0, str(bullet)


def _normalized_rows(rows, column_count):
    normalized = []
    for row in rows:
        values = [str(value) for value in row]
        if len(values) < column_count:
            values.extend([""] * (column_count - len(values)))
        normalized.append(values[:column_count])
    return normalized


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _blank_slide(prs, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, "background", theme)
    return slide


def _corner_accents(slide, theme):
    """Subtle decorative squares for title/closing slides."""
    for i, name in enumerate(("gold", "teal", "accent")):
        sq = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    SLIDE_WIDTH - Inches(1.6) + Inches(0.42) * i,
                                    Inches(0.55), Inches(0.3), Inches(0.3))
        _fill(sq, theme, name)
        _line_none(sq)


def _strip_table_style(table):
    """Remove the default banded PowerPoint table style for a clean custom look."""
    tbl = table._tbl
    tblPr = tbl.tblPr
    if tblPr is not None:
        for child in list(tblPr):
            if child.tag == qn("a:tableStyleId"):
                tblPr.remove(child)


def _style_cell(cell, text, theme, fill_color, text_color, *, bold=False, size=11, align=None):
    cell.fill.solid()
    cell.fill.fore_color.rgb = _rgb(theme, fill_color)
    cell.margin_left = Inches(0.14)
    cell.margin_right = Inches(0.1)
    cell.margin_top = Inches(0.06)
    cell.margin_bottom = Inches(0.06)
    frame = cell.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align is not None:
        paragraph.alignment = align
    _style_paragraph(paragraph, theme, size, text_color, bold=bold)


def _add_footer(slide, deck_title, page, total, theme):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X, Inches(6.92), CONTENT_W, Inches(0.012))
    _fill(line, theme, "border")
    _line_none(line)
    _add_text(slide, deck_title, MARGIN_X, FOOTER_TOP, Inches(9.0), Inches(0.25), theme,
              "muted_text", int(theme["small_font_size"]))
    _add_text(slide, f"{page} / {total}", Inches(11.5), FOOTER_TOP, Inches(1.1), Inches(0.25),
              theme, "muted_text", int(theme["small_font_size"]), align=PP_ALIGN.RIGHT)


def _add_text(slide, text, x, y, width, height, theme, color, size, *,
              bold=False, align=None, anchor=None, spacing=None):
    box = slide.shapes.add_textbox(x, y, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    if anchor is not None:
        frame.vertical_anchor = anchor
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align is not None:
        paragraph.alignment = align
    _style_paragraph(paragraph, theme, size, color, bold=bold, spacing=spacing)


def _style_paragraph(paragraph, theme, size, color, *, bold=False, spacing=None):
    """Apply font sizing/colour plus latin + east-asian font names to every run."""
    latin = str(theme["font"])
    ea = str(theme.get("kr_font", theme["font"]))
    for run in paragraph.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(theme, color)
        run.font.name = latin
        _set_run_fonts(run, latin, ea)
        if spacing is not None:
            _set_letter_spacing(run, spacing)


def _set_run_fonts(run, latin, ea):
    rPr = run._r.get_or_add_rPr()
    for tag, face in (("a:latin", latin), ("a:ea", ea), ("a:cs", latin)):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", face)


def _set_letter_spacing(run, points):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(points * 100)))


def _set_background(slide, color_name, theme):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(theme, color_name)


def _fill(shape, theme, color_name):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(theme, color_name)
    shape.shadow.inherit = False


def _line(shape, theme, color_name, width=1.0):
    shape.line.color.rgb = _rgb(theme, color_name)
    shape.line.width = Pt(width)


def _line_none(shape):
    shape.line.fill.background()


def _rounded(shape, radius_in=0.1):
    """Set an absolute corner radius (in inches) on a rounded rectangle so small
    and large panels share a consistent, subtle curvature instead of capsules."""
    try:
        smaller = min(int(shape.width), int(shape.height))
        if smaller > 0:
            shape.adjustments[0] = min(0.5, float(Inches(radius_in)) / smaller)
    except (AttributeError, IndexError, ValueError, ZeroDivisionError):
        pass


def _soft_shadow(shape, *, blur=46000, dist=20000, direction=5400000,
                 alpha=30000, color="0B1B30"):
    """Apply a subtle Office-style outer drop shadow for card depth.

    alpha is shadow opacity in 1/1000 % (30000 = 30%); direction 5400000 = down.
    """
    spPr = shape._element.spPr
    for existing in spPr.findall(qn("a:effectLst")):
        spPr.remove(existing)
    spPr.append(parse_xml(
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:outerShdw blurRad="{blur}" dist="{dist}" dir="{direction}" rotWithShape="0">'
        f'<a:srgbClr val="{color}"><a:alpha val="{alpha}"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'
    ))


def _panel(slide, x, y, w, h, theme, fill_color, *, line_color="border",
           shadow=True, radius=0.1):
    """Create a rounded card panel: fill + optional hairline border + consistent
    radius + optional soft shadow. The primary building block for card visuals."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _fill(shape, theme, fill_color)
    if line_color:
        _line(shape, theme, line_color)
    else:
        _line_none(shape)
    _rounded(shape, radius)
    if shadow:
        _soft_shadow(shape)
    return shape


def _bar(slide, x, y, w, h, theme, color):
    """Create a flat (square-cornered, shadowless) accent bar or segment."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _fill(shape, theme, color)
    _line_none(shape)
    return shape


def _series_color(theme, index):
    series = theme.get("series") or ["accent"]
    return series[(index - 1) % len(series)] if index >= 1 else series[0]


def _rgb(theme, color_name):
    colors = theme.get("colors", {})
    value = str(colors.get(color_name, color_name)).lstrip("#")
    if len(value) != 6:
        raise SlideSpecError(f"Invalid color value for '{color_name}': {value}")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def parse_args():
    parser = argparse.ArgumentParser(description="Generate a readable PowerPoint deck from a JSON slide spec.")
    parser.add_argument("spec", type=Path, help="Path to the slide spec JSON file.")
    parser.add_argument("-o", "--output", type=Path, required=True, help="Output .pptx path.")
    return parser.parse_args()


def main():
    args = parse_args()
    try:
        output = build_pptx(args.spec, args.output)
    except (OSError, json.JSONDecodeError, SlideSpecError) as exc:
        raise SystemExit(f"slide_builder error: {exc}") from exc
    print(f"Created {output}")


if __name__ == "__main__":
    main()
