from __future__ import annotations

import sys
import tempfile
import unittest
from pathlib import Path

SKILL_DIR = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_DIR))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402

import pptx_helpers as H  # noqa: E402


class PptxHelpersTests(unittest.TestCase):
    def test_builds_and_reopens_deck(self):
        prs, blank = H.new_deck()
        blue = H.hexc("1F63D8")
        ink = H.hexc("#14203A")
        # slide 1
        s1 = H.add_slide(prs, blank, bg=H.hexc("FFFFFF"))
        H.box(s1, 0.7, 0.5, 3.0, 0.14, fill=blue, kind=H.MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        H.text(s1, "결론형 제목", 0.7, 0.7, 11.9, 1.0, 30, ink, bold=True)
        H.bullets(s1, ["근거 1", "근거 2"], 0.7, 1.9, 6.0, 2.0, 16, ink, marker_color=blue)
        H.chip(s1, 0.7, 4.2, 2.0, 0.4, "GA", H.hexc("2E9E6B"), H.hexc("FFFFFF"))
        # slide 2 with a mechanical grid table
        s2 = H.add_slide(prs, blank)
        H.grid_table(
            s2, 0.7, 1.5, [3.0, 3.0, 3.0], [0.6, 0.9],
            [[{"text": "A", "fill": H.hexc("0E1B33"), "color": H.hexc("FFFFFF"), "bold": True},
              {"text": "B", "fill": H.hexc("0E1B33"), "color": H.hexc("FFFFFF")},
              {"text": "C", "fill": H.hexc("0E1B33"), "color": H.hexc("FFFFFF")}],
             [{"text": "1"}, {"text": "2"}, {"text": "3"}]],
            line_color=H.hexc("D8DEE7"),
        )
        H.chevron(s2, 4.0, 5.0, 0.4, 0.4, blue)

        with tempfile.TemporaryDirectory() as d:
            out = Path(d) / "helpers_smoke.pptx"
            prs.save(out)
            self.assertTrue(out.is_file())
            reopened = Presentation(out)
            self.assertEqual(len(reopened.slides._sldIdLst), 2)

    def test_module_defines_no_baked_palette(self):
        """디자인 중립 보장: 모듈 레벨에 색(RGBColor) 상수가 없어야 한다."""
        baked = [name for name in dir(H)
                 if not name.startswith("_") and isinstance(getattr(H, name), RGBColor)]
        self.assertEqual(baked, [], f"pptx_helpers must not bake palette colors: {baked}")

    def test_default_font_is_overridable_string(self):
        self.assertIsInstance(H.DEFAULT_FONT, str)


if __name__ == "__main__":
    unittest.main()
