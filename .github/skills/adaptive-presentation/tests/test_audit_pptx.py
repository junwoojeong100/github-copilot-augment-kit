from __future__ import annotations

import argparse
import sys
import subprocess
import tempfile
import unittest
import zipfile
from argparse import Namespace
from pathlib import Path

SCRIPTS_DIR = Path(__file__).resolve().parents[1] / "scripts"
sys.path.insert(0, str(SCRIPTS_DIR))

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

import audit_pptx  # noqa: E402


def audit_args(deck: Path, **overrides) -> Namespace:
    values = {
        "deck": deck,
        "expected_slides": 1,
        "allow_bleed": set(),
        "allow_small_text": set(),
        "allow_overlap": set(),
        "allow_title_size": set(),
        "require_sources": set(),
        "bounds_tolerance": 0.02,
        "min_body_pt": 13.0,
        "min_title_pt": 26.0,
        "title_size_tolerance_pt": 0.5,
        "footer_top": 6.9,
        "min_small_text_chars": 10,
        "fail_small_text": True,
        "fail_unsized_runs": True,
        "fail_title_risks": True,
        "fail_title_consistency": True,
        "fail_overlaps": True,
        "json": None,
        "strict": True,
    }
    values.update(overrides)
    return Namespace(**values)


class AuditPptxTests(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(self.temp_dir.cleanup)
        self.work_dir = Path(self.temp_dir.name)

    def make_deck(self, body_size: float = 13.0, explicit_body_size: bool = True) -> Path:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(0.7))
        title.text = "Conclusion title"
        title.text_frame.paragraphs[0].runs[0].font.size = Pt(30)
        body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        body.text = "This body sentence is long enough to be classified."
        if explicit_body_size:
            body.text_frame.paragraphs[0].runs[0].font.size = Pt(body_size)
        slide.shapes.add_shape(1, Inches(1), Inches(4), Inches(2), Inches(1))
        out = self.work_dir / "deck.pptx"
        prs.save(out)
        return out

    def test_strict_typography_passes_at_thirteen_points(self):
        report, failures = audit_pptx.audit(audit_args(self.make_deck()))
        self.assertEqual(failures, [])
        self.assertEqual(report["empty_text_frames"], [])

    def test_required_source_must_be_in_footer(self):
        deck = self.make_deck()
        _, failures = audit_pptx.audit(
            audit_args(deck, require_sources={1})
        )
        self.assertTrue(any("source footer missing" in item for item in failures))

        prs = Presentation(deck)
        empty_source = prs.slides[0].shapes.add_textbox(
            Inches(1), Inches(6.95), Inches(10), Inches(0.2)
        )
        empty_source.text = "Source: "
        empty_source.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
        prs.save(deck)
        _, failures = audit_pptx.audit(
            audit_args(deck, require_sources={1})
        )
        self.assertTrue(any("source footer missing" in item for item in failures))

        prs = Presentation(deck)
        body_source = prs.slides[0].shapes.add_textbox(
            Inches(1), Inches(5), Inches(10), Inches(0.2)
        )
        body_source.text = "Source: https://example.com/not-a-footer"
        body_source.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
        prs.save(deck)
        _, failures = audit_pptx.audit(
            audit_args(deck, require_sources={1})
        )
        self.assertTrue(any("source footer missing" in item for item in failures))

        prs = Presentation(deck)
        for source_shape in list(prs.slides[0].shapes)[-2:]:
            source_shape._element.getparent().remove(source_shape._element)
        footer_source = prs.slides[0].shapes.add_textbox(
            Inches(1), Inches(6.95), Inches(10), Inches(0.2)
        )
        footer_source.text = "출처: https://example.com/canonical"
        footer_source.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
        prs.save(deck)
        report, failures = audit_pptx.audit(
            audit_args(deck, require_sources={1})
        )
        self.assertEqual(failures, [])
        self.assertEqual(report["slides_with_footer_sources"], [1])
        self.assertEqual(report["missing_required_source_slides"], [])

    def test_nonfinite_thresholds_are_rejected(self):
        with self.assertRaises(argparse.ArgumentTypeError):
            audit_pptx.positive_float("nan")
        with self.assertRaises(argparse.ArgumentTypeError):
            audit_pptx.nonnegative_float("inf")

    def test_small_text_exception_is_slide_scoped(self):
        deck = self.make_deck(body_size=12)
        _, failures = audit_pptx.audit(audit_args(deck))
        self.assertTrue(any("below 13.0 pt" in item for item in failures))
        report, failures = audit_pptx.audit(
            audit_args(deck, allow_small_text={1})
        )
        self.assertEqual(failures, [])
        self.assertTrue(report["small_text_body_candidates"][0]["allowed"])

    def test_fractional_size_below_threshold_fails_before_rounding(self):
        _, failures = audit_pptx.audit(
            audit_args(self.make_deck(body_size=12.99))
        )
        self.assertTrue(any("below 13.0 pt" in item for item in failures))

    def test_compact_secondary_annotation_is_reported_as_label(self):
        deck = self.make_deck()
        prs = Presentation(deck)
        annotation = prs.slides[0].shapes.add_textbox(
            Inches(1), Inches(3.3), Inches(8), Inches(0.3)
        )
        annotation.text = (
            "Compact secondary annotation may use a smaller role-specific size."
        )
        annotation.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        prs.save(deck)

        report, failures = audit_pptx.audit(audit_args(deck))
        self.assertEqual(failures, [])
        self.assertEqual(report["small_text_body_candidates"], [])
        self.assertEqual(len(report["small_text_label_candidates"]), 1)

    def test_split_runs_cannot_bypass_body_length_threshold(self):
        deck = self.make_deck()
        prs = Presentation(deck)
        paragraph = prs.slides[0].shapes[1].text_frame.paragraphs[0]
        paragraph.clear()
        for value in ("Split ", "body ", "across ", "short ", "runs."):
            run = paragraph.add_run()
            run.text = value
            run.font.size = Pt(12)
        prs.save(deck)

        report, failures = audit_pptx.audit(audit_args(deck))
        self.assertGreaterEqual(len(report["small_text_body_candidates"]), 1)
        self.assertTrue(any("below 13.0 pt" in item for item in failures))

    def test_unsized_nonempty_run_fails(self):
        _, failures = audit_pptx.audit(
            audit_args(self.make_deck(explicit_body_size=False))
        )
        self.assertTrue(any("no explicit font size" in item for item in failures))

    def test_cover_requires_title_sized_text(self):
        deck = self.make_deck()
        prs = Presentation(deck)
        prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.size = Pt(20)
        prs.save(deck)
        report, failures = audit_pptx.audit(audit_args(deck))
        self.assertEqual(report["title_risks"][0]["slide"], 1)
        self.assertTrue(any("no text at or above" in item for item in failures))

    def test_content_title_sizes_are_consistent_while_cover_is_excluded(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        cover = prs.slides.add_slide(blank)
        cover_title = cover.shapes.add_textbox(
            Inches(1), Inches(1.2), Inches(10), Inches(0.9)
        )
        cover_title.text = "Cover title"
        cover_title.text_frame.paragraphs[0].runs[0].font.size = Pt(40)

        for number, size in enumerate((32, 32, 34), 1):
            slide = prs.slides.add_slide(blank)
            title = slide.shapes.add_textbox(
                Inches(1), Inches(0.7), Inches(10), Inches(1.1)
            )
            title.text = f"Content title {number}"
            title.text_frame.paragraphs[0].runs[0].font.size = Pt(size)
            body = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(8), Inches(1)
            )
            body.text = "This body sentence is long enough to be classified."
            body.text_frame.paragraphs[0].runs[0].font.size = Pt(15)

        deck = self.work_dir / "title-sizes.pptx"
        prs.save(deck)
        report, failures = audit_pptx.audit(
            audit_args(deck, expected_slides=4)
        )
        self.assertEqual(report["content_title_reference_pt"], 32)
        self.assertEqual(
            [
                item["slide"]
                for item in report["unexpected_title_size_inconsistencies"]
            ],
            [4],
        )
        self.assertTrue(
            any("title(s) differ" in item for item in failures)
        )

        report, failures = audit_pptx.audit(
            audit_args(
                deck,
                expected_slides=4,
                allow_title_size={4},
            )
        )
        self.assertEqual(failures, [])
        self.assertTrue(report["title_size_inconsistencies"][0]["allowed"])

    def test_table_cell_text_is_audited(self):
        deck = self.make_deck()
        prs = Presentation(deck)
        slide = prs.slides[0]
        table = slide.shapes.add_table(
            1, 1, Inches(1), Inches(4), Inches(8), Inches(1)
        ).table
        run = table.cell(0, 0).text_frame.paragraphs[0].add_run()
        run.text = "Table body content that must meet the same size threshold."
        run.font.size = Pt(12)
        prs.save(deck)

        report, failures = audit_pptx.audit(audit_args(deck))
        candidates = report["small_text_body_candidates"]
        self.assertTrue(any("cell 1,1" in item["shape"] for item in candidates))
        self.assertTrue(any("below 13.0 pt" in item for item in failures))

    def test_overlapping_text_frames_fail_and_allow_slide_exception(self):
        deck = self.make_deck()
        prs = Presentation(deck)
        overlap = prs.slides[0].shapes.add_textbox(
            Inches(1.5), Inches(2.2), Inches(4), Inches(0.6)
        )
        overlap.text = "A second body sentence overlaps the first text frame."
        overlap.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        prs.save(deck)

        report, failures = audit_pptx.audit(audit_args(deck))
        self.assertTrue(
            any(
                item["kind"] == "text_text"
                for item in report["unexpected_overlap_candidates"]
            )
        )
        self.assertTrue(any("geometry overlap" in item for item in failures))

        report, failures = audit_pptx.audit(
            audit_args(deck, allow_overlap={1})
        )
        self.assertEqual(failures, [])
        self.assertTrue(report["overlap_candidates"][0]["allowed"])

    def test_shape_and_text_container_overlaps_are_reported(self):
        deck = self.make_deck()
        prs = Presentation(deck)
        slide = prs.slides[0]
        chevron = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            Inches(2.9),
            Inches(4.3),
            Inches(0.2),
            Inches(0.3),
        )
        self.assertIsNotNone(chevron)
        text = slide.shapes.add_textbox(
            Inches(1.2), Inches(4.72), Inches(1.6), Inches(0.5)
        )
        text.text = "Text protruding below its intended card."
        text.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        prs.save(deck)

        report, failures = audit_pptx.audit(audit_args(deck))
        kinds = {
            item["kind"] for item in report["unexpected_overlap_candidates"]
        }
        self.assertIn("shape_shape", kinds)
        self.assertIn("text_container", kinds)
        self.assertTrue(any("geometry overlap" in item for item in failures))

    def test_text_table_cannot_use_decorative_bleed_exception(self):
        deck = self.make_deck()
        prs = Presentation(deck)
        table = prs.slides[0].shapes.add_table(
            1, 1, Inches(12.5), Inches(4), Inches(2), Inches(1)
        ).table
        table.cell(0, 0).text = "Visible text"
        prs.save(deck)

        report, failures = audit_pptx.audit(
            audit_args(deck, allow_bleed={1})
        )
        self.assertEqual(len(report["unexpected_out_of_bounds"]), 1)
        self.assertTrue(any("without allow-bleed" in item for item in failures))

    def test_json_report_cannot_alias_input_deck(self):
        deck = self.make_deck()
        alias = self.work_dir / "report.json"
        alias.hardlink_to(deck)
        result = subprocess.run(
            [
                sys.executable,
                "-B",
                str(SCRIPTS_DIR / "audit_pptx.py"),
                str(deck),
                "--json",
                str(alias),
            ],
            capture_output=True,
            text=True,
            check=False,
        )
        self.assertEqual(result.returncode, 2)
        self.assertTrue(zipfile.is_zipfile(deck))


if __name__ == "__main__":
    unittest.main()
