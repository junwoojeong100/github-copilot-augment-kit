from __future__ import annotations

import sys
import tempfile
import unittest
import zipfile
from pathlib import Path

from lxml import etree

SCRIPTS_DIR = Path(__file__).resolve().parents[1] / "scripts"
sys.path.insert(0, str(SCRIPTS_DIR))

from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

import audit_pptx  # noqa: E402


class RepairRiskTests(unittest.TestCase):
    def _deck_with_effectlst(self, count: int):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
        spPr = sp._element.spPr
        for _ in range(count):
            spPr.append(spPr.makeelement(qn("a:effectLst"), {}))
        return prs

    def test_flags_duplicate_effectlst(self):
        risks = audit_pptx.detect_repair_risks(self._deck_with_effectlst(2))
        self.assertEqual(len(risks), 1)
        self.assertEqual(risks[0]["slide"], 1)
        self.assertIn("effectLst", risks[0]["duplicate_children"])

    def test_single_effectlst_is_clean(self):
        self.assertEqual(audit_pptx.detect_repair_risks(self._deck_with_effectlst(1)), [])

    def test_no_effects_is_clean(self):
        self.assertEqual(audit_pptx.detect_repair_risks(self._deck_with_effectlst(0)), [])

    def test_package_scan_includes_slide_layout_parts(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            source = Path(temp_dir) / "source.pptx"
            broken = Path(temp_dir) / "broken.pptx"
            self._deck_with_effectlst(0).save(source)
            with zipfile.ZipFile(source) as src, zipfile.ZipFile(
                broken, "w", zipfile.ZIP_DEFLATED
            ) as dst:
                for name in src.namelist():
                    content = src.read(name)
                    if name == "ppt/slideLayouts/slideLayout1.xml":
                        root = etree.fromstring(content)
                        sp_pr = next(
                            element
                            for element in root.iter()
                            if etree.QName(element).localname == "spPr"
                        )
                        sp_pr.append(etree.Element(qn("a:effectLst")))
                        sp_pr.append(etree.Element(qn("a:effectLst")))
                        content = etree.tostring(
                            root, xml_declaration=True, encoding="UTF-8"
                        )
                    dst.writestr(name, content)

            risks = audit_pptx.detect_package_repair_risks(broken)
            layout_risks = [
                risk
                for risk in risks
                if risk["part"] == "ppt/slideLayouts/slideLayout1.xml"
            ]
            self.assertEqual(len(layout_risks), 1)
            self.assertIsNone(layout_risks[0]["slide"])
            self.assertIn("effectLst", layout_risks[0]["duplicate_children"])

    def test_package_scan_ignores_comments_and_custom_sppr_namespaces(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            source = Path(temp_dir) / "source.pptx"
            augmented = Path(temp_dir) / "augmented.pptx"
            self._deck_with_effectlst(0).save(source)
            with zipfile.ZipFile(source) as src, zipfile.ZipFile(
                augmented, "w", zipfile.ZIP_DEFLATED
            ) as dst:
                for name in src.namelist():
                    content = src.read(name)
                    if name == "ppt/slides/slide1.xml":
                        root = etree.fromstring(content)
                        sp_pr = next(
                            element
                            for element in root.iter()
                            if element.tag == qn("p:spPr")
                        )
                        sp_pr.append(etree.Comment("valid comment"))
                        content = etree.tostring(
                            root, xml_declaration=True, encoding="UTF-8"
                        )
                    dst.writestr(name, content)
                dst.writestr(
                    "customXml/fake.xml",
                    (
                        '<root xmlns:x="urn:custom"><x:spPr>'
                        "<x:fill/><x:fill/></x:spPr></root>"
                    ),
                )

            self.assertEqual(
                audit_pptx.detect_package_repair_risks(augmented), []
            )


if __name__ == "__main__":
    unittest.main()
