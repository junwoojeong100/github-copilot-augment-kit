"""Design-neutral python-pptx convenience wrappers.

이 모듈은 **디자인 결정을 하지 않는다.** 색·크기·위치·글꼴은 모두 호출자가 인자로 넘긴다.
팔레트·테마·레이아웃 패밀리·슬라이드 유형(cover/statement/…)·자동 배치를 정의하지 않는다.
목적은 오직 python-pptx 보일러플레이트(텍스트 프레임 설정, 라운드 사각형, 그림자 XML, 그리드 표
계산 등)를 줄여 매 덱의 '자유로운' 시각 구성을 더 빨리 코딩하도록 돕는 것이다.

규칙:
- 여기에 기본 색/팔레트를 추가하지 말 것. 색은 항상 인자.
- 여기에 'cover/architecture' 같은 의미 컴포넌트를 추가하지 말 것. 그건 각 덱의 build 스크립트가 자유 구성.
- 유일한 편의 기본값은 폰트 폴백 문자열(override 가능)과 표준 16:9 치수뿐이다.
"""
from __future__ import annotations

from typing import Sequence

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

__all__ = [
    "PP_ALIGN", "MSO_ANCHOR", "MSO_AUTO_SIZE", "MSO_SHAPE", "RGBColor", "Inches", "Pt",
    "hexc", "new_deck", "add_slide", "soft_shadow", "box", "text",
    "bullets", "chip", "hline", "vline", "chevron", "grid_table",
]

DEFAULT_FONT = "Apple SD Gothic Neo"  # 폴백일 뿐, 호출자가 언제든 override


def hexc(value: str) -> RGBColor:
    """'1F63D8' 또는 '#1F63D8' → RGBColor."""
    return RGBColor.from_string(value.lstrip("#").upper())


def new_deck(width_in: float = 13.333, height_in: float = 7.5):
    """(prs, blank_layout) 반환. 기본 16:9는 치수일 뿐 디자인이 아니다."""
    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)
    return prs, prs.slide_layouts[6]


def add_slide(prs, blank_layout, bg: RGBColor | None = None):
    s = prs.slides.add_slide(blank_layout)
    if bg is not None:
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = bg
    return s


def soft_shadow(shape, color: str = "9AA6B8", blur: int = 90000, dist: int = 38000, alpha: int = 38000):
    """도형에 부드러운 외곽 그림자(옵션). 색/강도는 인자.

    기존 effectLst(예: python-pptx의 shadow.inherit=False가 만든 빈 요소)를 **재사용**해
    한 spPr 안에 effectLst가 둘 이상 생기지 않게 한다. effectLst가 중복되면 OOXML 스키마 위반이라
    PowerPoint가 파일 열 때 복구(repair)를 요구한다(LibreOffice·unzip은 관대해 통과).
    """
    spPr = shape._element.spPr
    effect = spPr.find(qn("a:effectLst"))
    if effect is None:
        effect = spPr.makeelement(qn("a:effectLst"), {})
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            ln.addnext(effect)  # effectLst는 스키마상 a:ln 다음에 온다
        else:
            spPr.append(effect)
    else:
        for child in list(effect):
            effect.remove(child)
    shdw = spPr.makeelement(qn("a:outerShdw"),
                            {"blurRad": str(blur), "dist": str(dist), "dir": "5400000", "rotWithShape": "0"})
    clr = spPr.makeelement(qn("a:srgbClr"), {"val": color.lstrip("#").upper()})
    a = spPr.makeelement(qn("a:alpha"), {"val": str(alpha)})
    clr.append(a)
    shdw.append(clr)
    effect.append(shdw)
    return shape


def box(slide, x, y, w, h, fill: RGBColor | None = None, *, line: RGBColor | None = None,
        line_w: float = 1.0, kind=MSO_SHAPE.RECTANGLE, radius: float | None = None, shadow=False):
    """도형(기본 사각형). fill은 위치/키워드 모두 허용. fill/line=None이면 각각 없음."""
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        soft_shadow(sp)
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def _apply_runs(tf, lines, size, color, bold, italic, font, align, spacing):
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color


def text(slide, s, x, y, w, h, size, color: RGBColor, *, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font=DEFAULT_FONT, spacing=1.05, italic=False, margin=0.02):
    """word_wrap 텍스트 상자. '\\n'은 문단 분리. 모든 run에 font/size/color 명시."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Inches(margin))
    _apply_runs(tf, str(s).split("\n"), size, color, bold, italic, font, align, spacing)
    return tb


def bullets(slide, items: Sequence[str], x, y, w, h, size, color: RGBColor, *, marker="•",
            marker_color: RGBColor | None = None, gap=6, spacing=1.12, font=DEFAULT_FONT, bold=False):
    """불릿 목록(문단마다 marker). marker='' 이면 marker 없음."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Inches(0.02))
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = spacing
        p.space_after = Pt(gap)
        if marker:
            rm = p.add_run(); rm.text = marker + "  "
            rm.font.size = Pt(size); rm.font.name = font; rm.font.bold = True
            rm.font.color.rgb = marker_color or color
        r = p.add_run(); r.text = str(it)
        r.font.size = Pt(size); r.font.name = font; r.font.color.rgb = color; r.font.bold = bold
    return tb


def chip(slide, x, y, w, h, s, fill: RGBColor, text_color: RGBColor, *, size=11,
         line: RGBColor | None = None, font=DEFAULT_FONT, bold=True):
    """알약형 라벨(가운데 정렬 텍스트)."""
    sp = box(slide, x, y, w, h, fill=fill, line=line, line_w=1.0, kind=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Inches(0.02))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(s)
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = font; r.font.color.rgb = text_color
    return sp


def hline(slide, x, y, w, color: RGBColor, weight=0.02):
    return box(slide, x, y, w, weight, fill=color)


def vline(slide, x, y, h, color: RGBColor, weight=0.04):
    return box(slide, x, y, weight, h, fill=color)


def chevron(slide, x, y, w, h, color: RGBColor):
    """흐름 화살표(오른쪽 방향)."""
    return box(slide, x, y, w, h, fill=color, kind=MSO_SHAPE.CHEVRON)


def grid_table(slide, x, y, col_widths: Sequence[float], row_heights: Sequence[float],
               cells, *, line_color: RGBColor | None = None, line_w=0.75, font=DEFAULT_FONT):
    """기계적 그리드 표. 스타일은 전부 셀 dict 인자로.

    cells: 2D 리스트(행×열). 각 셀은 None 또는 dict:
      {text, fill, color, size, bold, align, anchor, spacing, pad}
    반환: 그린 셀 도형 dict {(r,c): shape}.
    """
    shapes: dict = {}
    yy = y
    for r, rh in enumerate(row_heights):
        xx = x
        for c, cw in enumerate(col_widths):
            spec = cells[r][c] if (r < len(cells) and c < len(cells[r])) else None
            if spec is None:
                spec = {}
            fill = spec.get("fill")
            sp = box(slide, xx, yy, cw, rh, fill=fill, line=line_color, line_w=line_w)
            shapes[(r, c)] = sp
            txt = spec.get("text")
            if txt is not None and str(txt) != "":
                pad = spec.get("pad", 0.12)
                text(slide, txt, xx + pad, yy, cw - (2 * pad), rh, spec.get("size", 12.5),
                     spec.get("color", hexc("14203A")), bold=spec.get("bold", False),
                     align=spec.get("align", PP_ALIGN.LEFT), anchor=spec.get("anchor", MSO_ANCHOR.MIDDLE),
                     font=font, spacing=spec.get("spacing", 1.05))
            xx += cw
        yy += rh
    return shapes
