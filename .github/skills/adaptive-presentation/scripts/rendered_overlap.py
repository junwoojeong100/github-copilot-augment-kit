#!/usr/bin/env python3
"""Detect rendered text spill and collisions in a PPTX-derived PDF."""

from __future__ import annotations

import math
from pathlib import Path

import fitz
from pptx import Presentation


def normalize_text(text: str) -> str:
    return "".join(character.casefold() for character in text if character.isalnum())


def distance_to_rect(rect: fitz.Rect, x: float, y: float) -> float:
    dx = max(rect.x0 - x, 0.0, x - rect.x1)
    dy = max(rect.y0 - y, 0.0, y - rect.y1)
    return math.hypot(dx, dy)


def shape_records(prs: Presentation, slide, page: fitz.Page) -> list[dict]:
    x_scale = page.rect.width / prs.slide_width
    y_scale = page.rect.height / prs.slide_height
    records: list[dict] = []
    for z_index, shape in enumerate(slide.shapes):
        if getattr(shape, "has_table", False):
            text = " ".join(
                cell.text.strip()
                for row in shape.table.rows
                for cell in row.cells
                if cell.text.strip()
            )
        elif getattr(shape, "has_text_frame", False):
            text = " ".join(shape.text.split())
        else:
            continue
        normalized = normalize_text(text)
        if not normalized:
            continue
        records.append(
            {
                "z_index": z_index,
                "shape": getattr(shape, "name", f"shape-{shape.shape_id}"),
                "shape_id": shape.shape_id,
                "text": text,
                "normalized": normalized,
                "rect": fitz.Rect(
                    page.rect.x0 + shape.left * x_scale,
                    page.rect.y0 + shape.top * y_scale,
                    page.rect.x0 + (shape.left + shape.width) * x_scale,
                    page.rect.y0 + (shape.top + shape.height) * y_scale,
                ),
            }
        )
    return records


def span_records(page: fitz.Page) -> list[dict]:
    records: list[dict] = []
    for block in page.get_text("dict")["blocks"]:
        if block.get("type") != 0:
            continue
        for line in block["lines"]:
            for span in line["spans"]:
                text = " ".join(span["text"].split())
                normalized = normalize_text(text)
                if not normalized:
                    continue
                records.append(
                    {
                        "text": text,
                        "normalized": normalized,
                        "size_pt": round(float(span["size"]), 2),
                        "rect": fitz.Rect(span["bbox"]),
                    }
                )
    return records


def assign_spans_to_shapes(spans: list[dict], shapes: list[dict]) -> tuple[list[dict], int]:
    assigned: list[dict] = []
    unmapped = 0
    for span in spans:
        rect = span["rect"]
        center_x = (rect.x0 + rect.x1) / 2
        center_y = (rect.y0 + rect.y1) / 2
        candidates: list[tuple[float, int]] = []
        for shape_index, shape in enumerate(shapes):
            span_text = span["normalized"]
            shape_text = shape["normalized"]
            if span_text not in shape_text and shape_text not in span_text:
                continue
            frame = shape["rect"]
            frame_center_x = (frame.x0 + frame.x1) / 2
            frame_center_y = (frame.y0 + frame.y1) / 2
            score = distance_to_rect(frame, center_x, center_y) * 10
            score += math.hypot(
                center_x - frame_center_x,
                center_y - frame_center_y,
            ) / max(len(span_text), 1)
            candidates.append((score, shape_index))
        if not candidates:
            unmapped += 1
            continue
        record = dict(span)
        record["shape_index"] = min(candidates)[1]
        assigned.append(record)
    return assigned, unmapped


def rect_values(rect: fitz.Rect) -> list[float]:
    return [round(value, 2) for value in (rect.x0, rect.y0, rect.x1, rect.y1)]


def detect_span_overlaps(
    assigned_spans: list[dict],
    shapes: list[dict],
    *,
    min_ratio: float = 0.03,
    min_dimension_pt: float = 1.0,
) -> list[dict]:
    overlaps: dict[tuple[int, int], dict] = {}
    for index, first in enumerate(assigned_spans):
        for second in assigned_spans[index + 1 :]:
            first_shape = int(first["shape_index"])
            second_shape = int(second["shape_index"])
            if first_shape == second_shape:
                continue
            intersection = first["rect"] & second["rect"]
            if (
                intersection.is_empty
                or intersection.width < min_dimension_pt
                or intersection.height < min_dimension_pt
            ):
                continue
            smaller_area = min(first["rect"].get_area(), second["rect"].get_area())
            ratio = intersection.get_area() / max(smaller_area, 1.0)
            if ratio < min_ratio:
                continue
            key = tuple(sorted((first_shape, second_shape)))
            item = {
                "shape_a": shapes[first_shape]["shape"],
                "shape_a_text": shapes[first_shape]["text"][:160],
                "span_a": first["text"][:120],
                "shape_b": shapes[second_shape]["shape"],
                "shape_b_text": shapes[second_shape]["text"][:160],
                "span_b": second["text"][:120],
                "overlap_ratio": round(ratio, 3),
                "overlap_rect_pt": rect_values(intersection),
            }
            previous = overlaps.get(key)
            if previous is None or item["overlap_ratio"] > previous["overlap_ratio"]:
                overlaps[key] = item
    return list(overlaps.values())


def detect_span_overflow(
    assigned_spans: list[dict],
    shapes: list[dict],
    *,
    tolerance_pt: float = 4.0,
) -> list[dict]:
    by_shape: dict[int, list[dict]] = {}
    for span in assigned_spans:
        by_shape.setdefault(int(span["shape_index"]), []).append(span)

    findings: list[dict] = []
    for shape_index, spans in by_shape.items():
        rendered = fitz.Rect(spans[0]["rect"])
        for span in spans[1:]:
            rendered |= span["rect"]
        frame = shapes[shape_index]["rect"]
        edges = {
            "left": max(0.0, frame.x0 - rendered.x0),
            "top": max(0.0, frame.y0 - rendered.y0),
            "right": max(0.0, rendered.x1 - frame.x1),
            "bottom": max(0.0, rendered.y1 - frame.y1),
        }
        worst = max(edges.values())
        if worst <= tolerance_pt:
            continue
        findings.append(
            {
                "shape": shapes[shape_index]["shape"],
                "text": shapes[shape_index]["text"][:160],
                "max_overflow_pt": round(worst, 2),
                "overflow_edges_pt": {
                    edge: round(value, 2)
                    for edge, value in edges.items()
                    if value > tolerance_pt
                },
                "frame_rect_pt": rect_values(frame),
                "rendered_rect_pt": rect_values(rendered),
            }
        )
    return findings


def audit_rendered_text(
    deck: Path,
    pdf: Path,
    *,
    allowed_slides: set[int] | None = None,
    overflow_tolerance_pt: float = 4.0,
) -> dict:
    allowed = allowed_slides or set()
    prs = Presentation(deck)
    overlap_findings: list[dict] = []
    overflow_findings: list[dict] = []
    unmapped_spans = 0

    with fitz.open(pdf) as document:
        if len(document) != len(prs.slides):
            raise ValueError(
                "Rendered PDF slide count differs from PPTX while checking overlaps: "
                f"{len(document)} vs {len(prs.slides)}"
            )
        for slide_number, (slide, page) in enumerate(
            zip(prs.slides, document), 1
        ):
            shapes = shape_records(prs, slide, page)
            assigned, unmapped = assign_spans_to_shapes(span_records(page), shapes)
            unmapped_spans += unmapped
            for finding in detect_span_overlaps(assigned, shapes):
                finding["slide"] = slide_number
                finding["allowed"] = slide_number in allowed
                overlap_findings.append(finding)
            for finding in detect_span_overflow(
                assigned,
                shapes,
                tolerance_pt=overflow_tolerance_pt,
            ):
                finding["slide"] = slide_number
                overflow_findings.append(finding)

    unexpected = [
        finding for finding in overlap_findings if not finding["allowed"]
    ]
    return {
        "rendered_text_overlaps": overlap_findings,
        "unexpected_rendered_text_overlaps": unexpected,
        "rendered_text_overflow_candidates": overflow_findings,
        "unmapped_rendered_text_spans": unmapped_spans,
    }
