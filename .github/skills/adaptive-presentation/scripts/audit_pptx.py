#!/usr/bin/env python3
"""Audit a PowerPoint deck for structural and typography risks."""

from __future__ import annotations

import argparse
import json
import math
import re
import sys
import zipfile
from collections import Counter
from pathlib import Path
from statistics import median

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches
from lxml import etree

from tooling import paths_collide

SHAPE_PROPERTY_TAGS = {
    qn("p:spPr"),
    qn("p:grpSpPr"),
    qn("a:spPr"),
    qn("a:grpSpPr"),
    qn("c:spPr"),
    "{http://schemas.openxmlformats.org/drawingml/2006/diagram}spPr",
}


def positive_int(value: str) -> int:
    parsed = int(value)
    if parsed <= 0:
        raise argparse.ArgumentTypeError("value must be greater than zero")
    return parsed


def positive_float(value: str) -> float:
    parsed = float(value)
    if not math.isfinite(parsed) or parsed <= 0:
        raise argparse.ArgumentTypeError("value must be greater than zero")
    return parsed


def nonnegative_float(value: str) -> float:
    parsed = float(value)
    if not math.isfinite(parsed) or parsed < 0:
        raise argparse.ArgumentTypeError("value must be non-negative")
    return parsed


def parse_slide_set(value: str) -> set[int]:
    if not value.strip():
        return set()
    slides: set[int] = set()
    for item in value.split(","):
        item = item.strip()
        if not item:
            continue
        if "-" in item:
            start_text, end_text = item.split("-", 1)
            start, end = int(start_text), int(end_text)
            if start > end:
                raise argparse.ArgumentTypeError(f"Invalid slide range: {item}")
            slides.update(range(start, end + 1))
        else:
            slides.add(int(item))
    if any(slide < 1 for slide in slides):
        raise argparse.ArgumentTypeError("Slide numbers must be positive")
    return slides


def shape_name(shape) -> str:
    return getattr(shape, "name", f"shape-{shape.shape_id}")


def iter_text_frames(
    shapes, slide_number: int, group_shapes: list[dict], *, in_group: bool = False
):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_shapes.append(
                {
                    "slide": slide_number,
                    "shape": shape_name(shape),
                    "child_count": len(shape.shapes),
                    "note": "Child text is audited; child bounds require visual verification.",
                }
            )
            yield from iter_text_frames(
                shape.shapes, slide_number, group_shapes, in_group=True
            )
        elif getattr(shape, "has_table", False):
            for row_number, row in enumerate(shape.table.rows, 1):
                for column_number, cell in enumerate(row.cells, 1):
                    yield (
                        cell.text_frame,
                        None if in_group else shape.top,
                        None if in_group else shape.height,
                        f"{shape_name(shape)} cell {row_number},{column_number}",
                        False,
                    )
        elif getattr(shape, "has_text_frame", False):
            yield (
                shape.text_frame,
                None if in_group else shape.top,
                None if in_group else shape.height,
                shape_name(shape),
                shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX,
            )


def shape_has_visible_text(shape) -> bool:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return any(shape_has_visible_text(child) for child in shape.shapes)
    if getattr(shape, "has_table", False):
        return any(
            cell.text.strip()
            for row in shape.table.rows
            for cell in row.cells
        )
    if getattr(shape, "has_chart", False):
        return True
    return getattr(shape, "has_text_frame", False) and bool(shape.text.strip())


def shape_text_content(shape) -> str:
    if getattr(shape, "has_table", False):
        return " ".join(
            cell.text.strip()
            for row in shape.table.rows
            for cell in row.cells
            if cell.text.strip()
        )
    if getattr(shape, "has_text_frame", False):
        return shape.text.strip()
    return ""


def looks_like_label(text: str) -> bool:
    compact = " ".join(text.split())
    if not compact or "\n" in text or len(compact) > 28:
        return False
    latin_letters = [char for char in compact if char.isalpha() and char.isascii()]
    if latin_letters and compact.upper() == compact:
        return True
    return len(compact.split()) <= 3 and len(compact) <= 20


def has_source_citation(text: str) -> bool:
    for line in text.splitlines():
        stripped = line.strip()
        lowered = stripped.lower()
        for marker in ("source:", "출처:"):
            if lowered.startswith(marker):
                citation = stripped[len(marker) :].strip()
                if len(citation) >= 3:
                    return True
    return False


def shape_bounds(shape) -> tuple[int, int, int, int]:
    return (
        shape.left,
        shape.top,
        shape.left + shape.width,
        shape.top + shape.height,
    )


def bounds_area(bounds: tuple[int, int, int, int]) -> int:
    return max(0, bounds[2] - bounds[0]) * max(0, bounds[3] - bounds[1])


def intersection_bounds(
    first: tuple[int, int, int, int],
    second: tuple[int, int, int, int],
    tolerance: int,
) -> tuple[int, int, int, int] | None:
    overlap = (
        max(first[0], second[0]),
        max(first[1], second[1]),
        min(first[2], second[2]),
        min(first[3], second[3]),
    )
    if (
        overlap[2] - overlap[0] <= tolerance
        or overlap[3] - overlap[1] <= tolerance
    ):
        return None
    return overlap


def detect_geometry_overlap_candidates(
    prs: Presentation,
    tolerance: int,
    allowed_slides: set[int],
) -> list[dict]:
    findings: list[dict] = []
    thin_threshold = max(tolerance, Inches(0.08))

    for slide_number, slide in enumerate(prs.slides, 1):
        shapes = [
            (z_index, shape, shape_bounds(shape))
            for z_index, shape in enumerate(slide.shapes)
            if shape.shape_type != MSO_SHAPE_TYPE.GROUP
        ]
        text_shapes = [
            item for item in shapes if shape_text_content(item[1])
        ]

        for index, (_, first, first_bounds) in enumerate(text_shapes):
            for _, second, second_bounds in text_shapes[index + 1 :]:
                overlap = intersection_bounds(
                    first_bounds, second_bounds, tolerance
                )
                if overlap is None:
                    continue
                overlap_area = bounds_area(overlap)
                smaller_area = min(
                    bounds_area(first_bounds), bounds_area(second_bounds)
                )
                findings.append(
                    {
                        "slide": slide_number,
                        "kind": "text_text",
                        "shape_a": shape_name(first),
                        "shape_a_text": shape_text_content(first)[:120],
                        "shape_b": shape_name(second),
                        "shape_b_text": shape_text_content(second)[:120],
                        "overlap_width_in": round(
                            (overlap[2] - overlap[0]) / Inches(1), 3
                        ),
                        "overlap_height_in": round(
                            (overlap[3] - overlap[1]) / Inches(1), 3
                        ),
                        "smaller_overlap_ratio": round(
                            overlap_area / max(smaller_area, 1), 3
                        ),
                        "allowed": slide_number in allowed_slides,
                    }
                )

        non_text_shapes = [
            item
            for item in shapes
            if not shape_has_visible_text(item[1])
            and item[1].shape_type != MSO_SHAPE_TYPE.LINE
            and item[1].width > thin_threshold
            and item[1].height > thin_threshold
        ]
        for index, (_, first, first_bounds) in enumerate(non_text_shapes):
            for _, second, second_bounds in non_text_shapes[index + 1 :]:
                overlap = intersection_bounds(
                    first_bounds, second_bounds, tolerance
                )
                if overlap is None:
                    continue
                overlap_area = bounds_area(overlap)
                first_area = bounds_area(first_bounds)
                second_area = bounds_area(second_bounds)
                smaller_ratio = overlap_area / max(
                    min(first_area, second_area), 1
                )
                larger_ratio = overlap_area / max(
                    max(first_area, second_area), 1
                )
                if (
                    smaller_ratio < 0.10
                    or larger_ratio < 0.001
                    or smaller_ratio > 0.95
                ):
                    continue
                findings.append(
                    {
                        "slide": slide_number,
                        "kind": "shape_shape",
                        "shape_a": shape_name(first),
                        "shape_b": shape_name(second),
                        "overlap_width_in": round(
                            (overlap[2] - overlap[0]) / Inches(1), 3
                        ),
                        "overlap_height_in": round(
                            (overlap[3] - overlap[1]) / Inches(1), 3
                        ),
                        "smaller_overlap_ratio": round(smaller_ratio, 3),
                        "allowed": slide_number in allowed_slides,
                    }
                )

        for text_z, text_shape, text_bounds in text_shapes:
            if (
                text_shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX
                and not getattr(text_shape, "has_table", False)
            ):
                continue
            text_area = bounds_area(text_bounds)
            center_x = (text_bounds[0] + text_bounds[2]) / 2
            center_y = (text_bounds[1] + text_bounds[3]) / 2
            containers: list[tuple[int, object, tuple[int, int, int, int], float]] = []
            for shape_z, shape, bounds in non_text_shapes:
                if shape_z >= text_z:
                    continue
                if shape.width < text_shape.width or shape.height < text_shape.height:
                    continue
                if not (
                    bounds[0] <= center_x <= bounds[2]
                    and bounds[1] <= center_y <= bounds[3]
                ):
                    continue
                overlap = intersection_bounds(text_bounds, bounds, 0)
                if overlap is None:
                    continue
                coverage = bounds_area(overlap) / max(text_area, 1)
                if 0.50 <= coverage < 0.98:
                    containers.append((shape_z, shape, bounds, coverage))
            if not containers:
                continue
            _, container, container_bounds, coverage = min(
                containers,
                key=lambda item: bounds_area(item[2]),
            )
            overflow = {
                "left": max(0, container_bounds[0] - text_bounds[0]),
                "top": max(0, container_bounds[1] - text_bounds[1]),
                "right": max(0, text_bounds[2] - container_bounds[2]),
                "bottom": max(0, text_bounds[3] - container_bounds[3]),
            }
            if max(overflow.values()) <= tolerance:
                continue
            findings.append(
                {
                    "slide": slide_number,
                    "kind": "text_container",
                    "shape_a": shape_name(text_shape),
                    "shape_a_text": shape_text_content(text_shape)[:120],
                    "shape_b": shape_name(container),
                    "text_coverage_ratio": round(coverage, 3),
                    "overflow_in": {
                        edge: round(value / Inches(1), 3)
                        for edge, value in overflow.items()
                        if value > tolerance
                    },
                    "allowed": slide_number in allowed_slides,
                }
            )

    return findings


def detect_title_size_consistency(
    prs: Presentation,
    *,
    min_title_pt: float,
    tolerance_pt: float,
    allowed_slides: set[int],
) -> dict:
    title_rows: list[dict] = []
    for slide_number, slide in enumerate(prs.slides, 1):
        candidates: list[dict] = []
        for shape in slide.shapes:
            if (
                not getattr(shape, "has_text_frame", False)
                or not shape.text.strip()
                or shape.top < Inches(0.25)
                or shape.top > Inches(1.85)
            ):
                continue
            sizes = sorted(
                {
                    round(run.font.size.pt, 2)
                    for paragraph in shape.text_frame.paragraphs
                    for run in paragraph.runs
                    if run.text.strip() and run.font.size
                }
            )
            if not sizes or max(sizes) < min_title_pt:
                continue
            candidates.append(
                {
                    "slide": slide_number,
                    "shape": shape_name(shape),
                    "text": " ".join(shape.text.split())[:160],
                    "top_in": round(shape.top / Inches(1), 3),
                    "left_in": round(shape.left / Inches(1), 3),
                    "width_in": round(shape.width / Inches(1), 3),
                    "size_pt": max(sizes),
                    "run_sizes_pt": sizes,
                }
            )
        if candidates:
            title_rows.append(
                max(
                    candidates,
                    key=lambda item: (
                        item["size_pt"],
                        item["width_in"],
                        -item["top_in"],
                    ),
                )
            )

    if len(title_rows) < 2:
        return {
            "title_rows": title_rows,
            "content_title_reference_pt": None,
            "content_title_size_range_pt": 0.0,
            "title_size_inconsistencies": [],
            "unexpected_title_size_inconsistencies": [],
        }

    top_buckets = Counter(round(item["top_in"] * 10) for item in title_rows)
    dominant_top_bucket = min(
        top_buckets,
        key=lambda bucket: (-top_buckets[bucket], bucket),
    )
    content_titles = [
        item
        for item in title_rows
        if abs(item["top_in"] - dominant_top_bucket / 10) <= 0.11
    ]
    baseline_titles = [
        item
        for item in content_titles
        if item["slide"] not in allowed_slides
    ]
    if not baseline_titles:
        baseline_titles = content_titles
    reference = (
        round(median(item["size_pt"] for item in baseline_titles), 2)
        if baseline_titles
        else None
    )
    size_range = (
        round(
            max(item["size_pt"] for item in baseline_titles)
            - min(item["size_pt"] for item in baseline_titles),
            2,
        )
        if baseline_titles
        else 0.0
    )

    inconsistencies: list[dict] = []
    if reference is not None:
        for item in content_titles:
            delta = abs(item["size_pt"] - reference)
            mixed_runs = (
                max(item["run_sizes_pt"]) - min(item["run_sizes_pt"])
                > tolerance_pt
            )
            if delta <= tolerance_pt and not mixed_runs:
                continue
            issue = dict(item)
            issue.update(
                {
                    "reference_size_pt": reference,
                    "delta_pt": round(delta, 2),
                    "mixed_run_sizes": mixed_runs,
                    "allowed": item["slide"] in allowed_slides,
                }
            )
            inconsistencies.append(issue)

    unexpected = [
        item for item in inconsistencies if not item["allowed"]
    ]
    return {
        "title_rows": title_rows,
        "content_title_reference_pt": reference,
        "content_title_size_range_pt": size_range,
        "title_size_inconsistencies": inconsistencies,
        "unexpected_title_size_inconsistencies": unexpected,
    }


def detect_repair_risks(prs) -> list[dict]:
    """Find spPr elements with duplicate singleton children (invalid OOXML).

    In CT_ShapeProperties each child (xfrm, geometry, fill, ln, effectLst, ...) may appear
    at most once. e.g. two <a:effectLst> in one shape makes PowerPoint prompt to repair the
    file on open, even though LibreOffice and `unzip -t` accept it.
    """
    risks: list[dict] = []
    shape_property_tags = {qn("p:spPr"), qn("p:grpSpPr")}
    for number, slide in enumerate(prs.slides, 1):
        for spPr in (
            element
            for element in slide._element.iter()
            if element.tag in shape_property_tags
        ):
            counts = Counter(child.tag for child in spPr)
            dups = sorted(
                tag.split("}")[-1] for tag, count in counts.items() if count > 1
            )
            if dups:
                risks.append({"slide": number, "duplicate_children": dups})
    return risks


def detect_package_repair_risks(
    deck: Path, slide_part_order: dict[str, int] | None = None
) -> list[dict]:
    """Scan shape-property XML across slides, layouts, masters, and related parts."""
    risks: list[dict] = []
    with zipfile.ZipFile(deck) as archive:
        for part in sorted(
            name for name in archive.namelist() if name.endswith(".xml")
        ):
            root = etree.fromstring(archive.read(part))
            slide_match = re.fullmatch(r"ppt/slides/slide(\d+)\.xml", part)
            slide_number = (
                slide_part_order.get(part)
                if slide_part_order is not None
                else int(slide_match.group(1)) if slide_match else None
            )
            for element in root.iter():
                if not isinstance(element.tag, str):
                    continue
                if element.tag not in SHAPE_PROPERTY_TAGS:
                    continue
                counts = Counter(
                    child.tag
                    for child in element
                    if isinstance(child.tag, str)
                )
                duplicates = sorted(
                    etree.QName(tag).localname
                    for tag, count in counts.items()
                    if count > 1
                )
                if duplicates:
                    risks.append(
                        {
                            "slide": slide_number,
                            "part": part,
                            "duplicate_children": duplicates,
                        }
                    )
    return risks


def audit(args: argparse.Namespace) -> tuple[dict, list[str]]:
    deck = args.deck.expanduser().resolve()
    if not deck.is_file():
        raise FileNotFoundError(deck)

    with zipfile.ZipFile(deck) as archive:
        corrupt_member = archive.testzip()

    prs = Presentation(deck)
    slide_count = len(prs.slides)
    for option, slides in (
        ("--allow-bleed", args.allow_bleed),
        ("--allow-small-text", args.allow_small_text),
        ("--allow-overlap", args.allow_overlap),
        ("--allow-title-size", args.allow_title_size),
    ):
        invalid = sorted(slide for slide in slides if slide > slide_count)
        if invalid:
            raise ValueError(
                f"{option} slides out of range 1-{slide_count}: {invalid}"
            )
    slide_part_order = {
        str(slide.part.partname).lstrip("/"): number
        for number, slide in enumerate(prs.slides, 1)
    }
    repair_risks = detect_package_repair_risks(deck, slide_part_order)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    tolerance = Inches(args.bounds_tolerance)
    footer_top = Inches(args.footer_top)

    fonts: Counter[str] = Counter()
    font_sizes: Counter[float] = Counter()
    text_chars: list[int] = []
    out_of_bounds: list[dict] = []
    small_text_body: list[dict] = []
    small_text_labels: list[dict] = []
    empty_text_frames: list[dict] = []
    title_risks: list[dict] = []
    unsized_runs: list[dict] = []
    group_shapes: list[dict] = []
    slides_with_sources: list[int] = []
    slides_with_footer_sources: list[int] = []
    overlap_candidates = detect_geometry_overlap_candidates(
        prs, tolerance, args.allow_overlap
    )
    title_consistency = detect_title_size_consistency(
        prs,
        min_title_pt=args.min_title_pt,
        tolerance_pt=args.title_size_tolerance_pt,
        allowed_slides=args.allow_title_size,
    )

    for slide_number, slide in enumerate(prs.slides, 1):
        slide_chars = 0
        max_font_size = 0.0
        has_source = False
        has_footer_source = False

        for shape in slide.shapes:
            left = shape.left
            top = shape.top
            right = shape.left + shape.width
            bottom = shape.top + shape.height
            if (
                left < -tolerance
                or top < -tolerance
                or right > slide_width + tolerance
                or bottom > slide_height + tolerance
            ):
                visible_text = shape_has_visible_text(shape)
                issue = {
                    "slide": slide_number,
                    "shape": shape_name(shape),
                    "left_in": round(left / Inches(1), 3),
                    "top_in": round(top / Inches(1), 3),
                    "width_in": round(shape.width / Inches(1), 3),
                    "height_in": round(shape.height / Inches(1), 3),
                    "has_visible_text": visible_text,
                    "allowed_bleed": (
                        slide_number in args.allow_bleed and not visible_text
                    ),
                }
                out_of_bounds.append(issue)

        for text_frame, top, frame_height, target_name, report_empty in iter_text_frames(
            slide.shapes, slide_number, group_shapes
        ):
            text = text_frame.text.strip()
            slide_chars += len(text)
            if not text and report_empty:
                empty_text_frames.append(
                    {"slide": slide_number, "shape": target_name}
                )
            if has_source_citation(text):
                has_source = True
                if top is not None and top >= footer_top:
                    has_footer_source = True

            for paragraph in text_frame.paragraphs:
                paragraph_text = paragraph.text.strip()
                for run in paragraph.runs:
                    if run.font.name:
                        fonts[run.font.name] += 1
                    if not run.font.size:
                        if run.text.strip():
                            unsized_runs.append(
                                {
                                    "slide": slide_number,
                                    "shape": target_name,
                                    "text": run.text.strip()[:120],
                                }
                            )
                        continue

                    raw_size = run.font.size.pt
                    reported_size = round(raw_size, 1)
                    font_sizes[reported_size] += 1
                    max_font_size = max(max_font_size, raw_size)
                    is_footer = top is not None and top >= footer_top
                    run_text = run.text.strip()
                    if (
                        run_text
                        and raw_size < args.min_body_pt
                        and not is_footer
                        and len(paragraph_text) >= args.min_small_text_chars
                    ):
                        issue = {
                            "slide": slide_number,
                            "shape": target_name,
                            "size_pt": reported_size,
                            "size_pt_raw": raw_size,
                            "text": run_text[:120],
                            "allowed": slide_number in args.allow_small_text,
                        }
                        is_compact_annotation = (
                            frame_height is not None
                            and frame_height <= Inches(0.36)
                            and len(paragraph_text) <= 100
                        )
                        if looks_like_label(paragraph_text) or is_compact_annotation:
                            small_text_labels.append(issue)
                        else:
                            small_text_body.append(issue)

        text_chars.append(slide_chars)
        if has_source:
            slides_with_sources.append(slide_number)
        if has_footer_source:
            slides_with_footer_sources.append(slide_number)
        if max_font_size < args.min_title_pt:
            title_risks.append(
                {
                    "slide": slide_number,
                    "max_font_size_pt": max_font_size,
                    "reason": "No text at or above minimum title size",
                }
            )

    unexpected_bounds = [
        item for item in out_of_bounds if not item["allowed_bleed"]
    ]
    unexpected_small_text_body = [
        item for item in small_text_body if not item["allowed"]
    ]
    unexpected_overlap_candidates = [
        item for item in overlap_candidates if not item["allowed"]
    ]
    required_source_slides = sorted(
        getattr(args, "require_sources", set())
    )
    missing_required_source_slides = sorted(
        set(required_source_slides) - set(slides_with_footer_sources)
    )
    report = {
        "deck": str(deck),
        "slides": len(prs.slides),
        "size_inches": {
            "width": round(slide_width / Inches(1), 3),
            "height": round(slide_height / Inches(1), 3),
        },
        "zip_integrity": "ok" if corrupt_member is None else corrupt_member,
        "fonts": fonts.most_common(),
        "font_sizes_pt": sorted(font_sizes.items(), reverse=True),
        "text_chars_per_slide": {
            "min": min(text_chars) if text_chars else 0,
            "average": round(sum(text_chars) / len(text_chars), 1)
            if text_chars
            else 0,
            "max": max(text_chars) if text_chars else 0,
            "values": text_chars,
        },
        "slides_with_sources": slides_with_sources,
        "slides_with_footer_sources": slides_with_footer_sources,
        "required_source_slides": required_source_slides,
        "missing_required_source_slides": missing_required_source_slides,
        "out_of_bounds": out_of_bounds,
        "unexpected_out_of_bounds": unexpected_bounds,
        "small_text_body_candidates": small_text_body,
        "unexpected_small_text_body_candidates": unexpected_small_text_body,
        "small_text_label_candidates": small_text_labels,
        "empty_text_frames": empty_text_frames,
        "title_risks": title_risks,
        **title_consistency,
        "unsized_runs": unsized_runs,
        "group_shapes": group_shapes,
        "overlap_candidates": overlap_candidates,
        "unexpected_overlap_candidates": unexpected_overlap_candidates,
        "ooxml_repair_risks": repair_risks,
    }

    failures: list[str] = []
    if corrupt_member is not None:
        failures.append(f"Corrupt ZIP member: {corrupt_member}")
    if repair_risks:
        failures.append(
            f"{len(repair_risks)} shape(s) have duplicate OOXML children "
            "(PowerPoint will prompt to repair)"
        )
    if args.expected_slides is not None and len(prs.slides) != args.expected_slides:
        failures.append(
            f"Expected {args.expected_slides} slides, found {len(prs.slides)}"
        )
    if unexpected_bounds:
        failures.append(
            f"{len(unexpected_bounds)} shapes exceed slide bounds without allow-bleed"
        )
    if args.fail_small_text and unexpected_small_text_body:
        failures.append(
            f"{len(unexpected_small_text_body)} likely body-text runs are below "
            f"{args.min_body_pt} pt"
        )
    if args.fail_title_risks and title_risks:
        failures.append(
            f"{len(title_risks)} slides have no text at or above {args.min_title_pt} pt"
        )
    if (
        args.fail_title_consistency
        and title_consistency["unexpected_title_size_inconsistencies"]
    ):
        failures.append(
            f"{len(title_consistency['unexpected_title_size_inconsistencies'])} "
            "content slide title(s) differ from the deck title-size standard"
        )
    if args.fail_unsized_runs and unsized_runs:
        failures.append(
            f"{len(unsized_runs)} non-empty text runs have no explicit font size"
        )
    if args.fail_overlaps and unexpected_overlap_candidates:
        failures.append(
            f"{len(unexpected_overlap_candidates)} geometry overlap candidate(s) "
            "require repair or --allow-overlap review"
        )
    if missing_required_source_slides:
        failures.append(
            "Required source footer missing on slide(s): "
            + ", ".join(str(slide) for slide in missing_required_source_slides)
        )

    return report, failures


def print_report(report: dict, failures: list[str]) -> None:
    print(f"Deck: {report['deck']}")
    print(
        f"Slides: {report['slides']} | "
        f"Size: {report['size_inches']['width']} × "
        f"{report['size_inches']['height']} in"
    )
    print(f"ZIP integrity: {report['zip_integrity']}")
    print(f"Fonts: {report['fonts']}")
    print(f"Font sizes (pt): {report['font_sizes_pt']}")
    density = report["text_chars_per_slide"]
    print(
        "Text chars/slide: "
        f"min={density['min']} avg={density['average']} max={density['max']}"
    )
    print(
        "Sources: "
        f"{len(report['slides_with_footer_sources'])}/{report['slides']} footer, "
        f"{len(report['slides_with_sources'])}/{report['slides']} anywhere"
    )
    if report["required_source_slides"]:
        print(
            "Required source slides: "
            f"{report['required_source_slides']} | "
            f"missing={report['missing_required_source_slides']}"
        )
    print(
        "Bounds: "
        f"{len(report['out_of_bounds'])} total, "
        f"{len(report['unexpected_out_of_bounds'])} unexpected"
    )
    print(
        "Small-text candidates: "
        f"{len(report['small_text_body_candidates'])} body, "
        f"{len(report['unexpected_small_text_body_candidates'])} unapproved body, "
        f"{len(report['small_text_label_candidates'])} label"
    )
    print(f"Unsized runs: {len(report['unsized_runs'])}")
    print(f"Group shapes: {len(report['group_shapes'])}")
    print(
        "Geometry overlaps: "
        f"{len(report.get('overlap_candidates', []))} total, "
        f"{len(report.get('unexpected_overlap_candidates', []))} unapproved"
    )
    print(f"Title risks: {len(report['title_risks'])}")
    print(
        "Title-size consistency: "
        f"reference={report.get('content_title_reference_pt')} pt, "
        f"range={report.get('content_title_size_range_pt')} pt, "
        f"{len(report.get('unexpected_title_size_inconsistencies', []))} unapproved"
    )
    print(f"OOXML repair risks: {len(report.get('ooxml_repair_risks', []))}")

    for label, items in (
        ("Unexpected bounds", report["unexpected_out_of_bounds"]),
        ("Small body text", report["small_text_body_candidates"]),
        ("Unsized runs", report["unsized_runs"]),
        ("Group shapes", report["group_shapes"]),
        ("Geometry overlaps", report.get("overlap_candidates", [])),
        ("Title risks", report["title_risks"]),
        (
            "Title-size inconsistencies",
            report.get("title_size_inconsistencies", []),
        ),
        ("OOXML repair risks", report.get("ooxml_repair_risks", [])),
    ):
        if items:
            print(f"\n{label}:")
            for item in items[:20]:
                print(f"  {item}")
            if len(items) > 20:
                print(f"  ... {len(items) - 20} more")

    if failures:
        print("\nFAIL:")
        for failure in failures:
            print(f"  - {failure}")
    else:
        print("\nPASS: no configured structural failures")


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Audit PPTX slide count, bounds, fonts, and text-density risks."
    )
    parser.add_argument("deck", type=Path, help="Path to the .pptx file")
    parser.add_argument(
        "--expected-slides",
        type=positive_int,
        help="Fail when the slide count differs from this value",
    )
    parser.add_argument(
        "--allow-bleed",
        type=parse_slide_set,
        default=set(),
        metavar="SLIDES",
        help="Comma/range list of slides allowed to contain out-of-bounds art, e.g. 1,5-6",
    )
    parser.add_argument(
        "--bounds-tolerance",
        type=nonnegative_float,
        default=0.02,
        metavar="INCHES",
        help="Bounds tolerance in inches (default: 0.02)",
    )
    parser.add_argument(
        "--min-body-pt",
        type=positive_float,
        default=13.0,
        help="Flag likely non-footer body text below this size (default: 13)",
    )
    parser.add_argument(
        "--min-title-pt",
        type=positive_float,
        default=26.0,
        help="Flag slides with no text at or above this size (default: 26)",
    )
    parser.add_argument(
        "--footer-top",
        type=nonnegative_float,
        default=6.9,
        metavar="INCHES",
        help="Text at or below this Y coordinate is treated as footer (default: 6.9)",
    )
    parser.add_argument(
        "--min-small-text-chars",
        type=positive_int,
        default=10,
        help="Ignore very short text below this character count (default: 10)",
    )
    parser.add_argument(
        "--allow-small-text",
        type=parse_slide_set,
        default=set(),
        metavar="SLIDES",
        help="Reviewed slides allowed to contain sub-minimum body text, e.g. 4,8-9",
    )
    parser.add_argument(
        "--allow-overlap",
        type=parse_slide_set,
        default=set(),
        metavar="SLIDES",
        help="Reviewed slides allowed intentional geometry overlap, e.g. 4,8-9",
    )
    parser.add_argument(
        "--allow-title-size",
        type=parse_slide_set,
        default=set(),
        metavar="SLIDES",
        help="Reviewed slides allowed a different content-title size, e.g. 6,12",
    )
    parser.add_argument(
        "--require-sources",
        type=parse_slide_set,
        default=set(),
        metavar="SLIDES",
        help="Slides with factual claims that must contain a Source:/출처: footer",
    )
    parser.add_argument(
        "--title-size-tolerance-pt",
        type=nonnegative_float,
        default=0.5,
        help="Allowed content-title size variation in points (default: 0.5)",
    )
    parser.add_argument(
        "--fail-small-text",
        action="store_true",
        help="Return nonzero when small-text candidates exist",
    )
    parser.add_argument(
        "--fail-unsized-runs",
        action="store_true",
        help="Return nonzero when non-empty runs have no explicit font size",
    )
    parser.add_argument(
        "--fail-title-risks",
        action="store_true",
        help="Return nonzero when a slide lacks title-sized text",
    )
    parser.add_argument(
        "--fail-title-consistency",
        action="store_true",
        help="Return nonzero when content-slide title sizes are inconsistent",
    )
    parser.add_argument(
        "--fail-overlaps",
        action="store_true",
        help="Return nonzero when unapproved geometry overlap candidates exist",
    )
    parser.add_argument(
        "--json",
        type=Path,
        help="Write the full audit report as JSON",
    )
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Enable typography, title, and overlap failures",
    )
    return parser


def main() -> int:
    parser = build_parser()
    args = parser.parse_args()
    if args.strict:
        args.fail_small_text = True
        args.fail_unsized_runs = True
        args.fail_title_risks = True
        args.fail_title_consistency = True
        args.fail_overlaps = True
    report, failures = audit(args)
    print_report(report, failures)

    if args.json:
        args.json = args.json.expanduser().resolve()
        if paths_collide(args.json, args.deck):
            parser.error("--json must not overwrite or alias the input deck")
        args.json.parent.mkdir(parents=True, exist_ok=True)
        args.json.write_text(
            json.dumps(report, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        print(f"JSON report: {args.json}")

    return 1 if failures else 0


if __name__ == "__main__":
    sys.exit(main())
