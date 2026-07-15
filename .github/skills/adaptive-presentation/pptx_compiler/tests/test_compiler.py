from __future__ import annotations

import sys
import tempfile
import unittest
import zipfile
from pathlib import Path

from PIL import Image
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN


SKILL_ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(SKILL_ROOT))

from pptx_compiler import (  # noqa: E402
    BlueprintLibrary,
    ContentItem,
    DesignDNA,
    DesignValidationError,
    LayoutError,
    LayoutMetrics,
    Palette,
    Region,
    ShapeLanguage,
    SlideCompiler,
    SlideFrame,
    Source,
    TextOverflowError,
    TextSpan,
    Typography,
)


def technical_design() -> DesignDNA:
    return DesignDNA(
        name="Technical Forge",
        concept_words=("precise", "composable", "operational"),
        visual_metaphor="programmable forge",
        palette=Palette(
            canvas="F4F1EA",
            canvas_alt="0B1017",
            surface="FFFFFF",
            surface_alt="141D28",
            ink="111820",
            ink_inverse="F7FAFC",
            muted_ink="5D6B78",
            muted_inverse="AAB6C2",
            primary="E65F2B",
            secondary="1677C8",
            accent="7856D8",
            border="D6D9DC",
            border_inverse="2A3948",
            success="238A57",
            warning="C48316",
            danger="C43D55",
            preview="8A4EC7",
        ),
        typography=Typography(
            body_font="Apple SD Gothic Neo",
            display_font="Apple SD Gothic Neo",
            mono_font="Menlo",
        ),
        shapes=ShapeLanguage(
            corner_style="subtle",
            spacing="balanced",
            depth="flat",
        ),
    )


def editorial_design() -> DesignDNA:
    return DesignDNA(
        name="Warm Editorial",
        concept_words=("measured", "human", "clear"),
        visual_metaphor="annotated field journal",
        palette=Palette(
            canvas="F7F2E8",
            canvas_alt="2A1F1A",
            surface="FFFFFA",
            surface_alt="392C25",
            ink="241C18",
            ink_inverse="FFF9F1",
            muted_ink="6E625B",
            muted_inverse="CDBFB2",
            primary="B34A32",
            secondary="4F6D5A",
            accent="92734B",
            border="D8CCC0",
            border_inverse="57463B",
            success="3F7654",
            warning="A66A22",
            danger="A33B45",
            preview="76528A",
        ),
        typography=Typography(
            body_font="Apple SD Gothic Neo",
            display_font="Georgia",
            mono_font="Menlo",
            cover_pt=42,
            title_pt=30,
            body_pt=19,
        ),
        shapes=ShapeLanguage(
            corner_style="rounded",
            spacing="spacious",
            depth="border",
        ),
    )


def build_sample(design: DesignDNA, output: Path):
    compiler = SlideCompiler(
        design,
        title="Compiler Sample",
        subject="Same semantics, different Design DNA",
    )

    cover = compiler.begin_slide(
        SlideFrame(
            title="",
            background_role="canvas_alt",
            inverse=True,
            show_header=False,
            show_footer=False,
        ),
        family="cover",
        variant="split",
    )
    cover.text(
        "kicker",
        "REUSABLE SLIDE COMPILER",
        role="label",
        color="primary",
        bold=True,
    )
    cover.text(
        "headline",
        "The compiler reuses mechanics,\nnot visual identity.",
        role="cover",
        color="ink_inverse",
        background="canvas_alt",
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    cover.text(
        "subtitle",
        "Design DNA remains an explicit input for every deck.",
        role="body",
        color="muted_inverse",
        background="canvas_alt",
    )
    cover.box("visual", fill="surface_alt", line="primary")
    cover.text(
        cover.region("visual", inset=0.45),
        "SPEC\n+\nDESIGN DNA\n→\nPPTX",
        role="metric",
        color="ink_inverse",
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        bold=True,
    )

    comparison = compiler.begin_slide(
        SlideFrame(
            title="Reusable mechanics do not require a reusable look",
            section="PRINCIPLE",
            number=2,
            source=Source(
                "Internal",
                "Compiler contract",
                "2026-07-15",
                "https://example.com/compiler",
            ),
        ),
        family="comparison",
        variant="axis",
    )
    comparison.comparison(
        "REUSE",
        ["Geometry safety", "Text fitting", "Source handling", "Validation"],
        "REGENERATE",
        ["Palette", "Typography", "Metaphor", "Layout rhythm"],
    )

    process = compiler.begin_slide(
        SlideFrame(
            title="Semantic steps compile through different design systems",
            section="FLOW",
            number=3,
        ),
        family="process",
        variant="rail",
        slots=4,
    )
    process.process(
        [
            ContentItem("SPEC", "Conclusion, evidence, relationship", "secondary"),
            ContentItem("DNA", "Palette, type, shape language", "accent"),
            ContentItem("COMPILE", "Regions, primitives, fitting", "primary"),
            ContentItem("VERIFY", "Audit, render, repair", "success"),
        ]
    )

    code = compiler.begin_slide(
        SlideFrame(
            title="Build scripts import the compiler instead of recreating helpers",
            section="CODE",
            number=4,
        ),
        family="code",
        variant="split",
    )
    code.code(
        [
            "compiler = SlideCompiler(design)",
            "slide = compiler.begin_slide(frame, family='process')",
            "slide.process(steps)",
            "compiler.save(output, expected_slides=4)",
        ],
        [
            ContentItem("Explicit design", "No built-in palette or branded template.", "primary"),
            ContentItem("Semantic layout", "The deck still chooses its blueprint family.", "secondary"),
            ContentItem("Strict output", "Bounds and diversity are checked before save.", "success"),
        ],
    )

    return compiler.save(
        output,
        expected_slides=4,
        min_layout_families=4,
    )


class CompilerTests(unittest.TestCase):
    def test_palette_rejects_low_contrast(self):
        with self.assertRaises(DesignValidationError):
            Palette(
                canvas="FFFFFF",
                canvas_alt="FFFFFF",
                surface="FFFFFF",
                surface_alt="FFFFFF",
                ink="EEEEEE",
                ink_inverse="EEEEEE",
                muted_ink="777777",
                muted_inverse="777777",
                primary="000000",
                secondary="111111",
                accent="222222",
                border="CCCCCC",
                border_inverse="333333",
                success="008000",
                warning="AA6600",
                danger="AA0000",
                preview="660099",
            )

    def test_all_blueprint_regions_stay_in_bounds(self):
        design = technical_design()
        library = BlueprintLibrary(design.layout, spacing=design.shapes.spacing)
        cases = (
            ("cover", "split", 4),
            ("cover", "centered", 4),
            ("statement", "hero-left", 4),
            ("statement", "centered", 4),
            ("comparison", "balanced", 4),
            ("comparison", "axis", 4),
            ("process", "rail", 4),
            ("process", "stair", 4),
            ("layered", "stack", 4),
            ("layered", "pyramid", 4),
            ("matrix", "table", 4),
            ("matrix", "quadrants", 4),
            ("architecture", "pipeline", 4),
            ("architecture", "hub", 4),
            ("code", "split", 4),
            ("code", "full", 4),
            ("roadmap", "rail", 4),
            ("roadmap", "gates", 4),
            ("custom", "free", 4),
        )
        for family, variant, slots in cases:
            with self.subTest(family=family, variant=variant):
                plan = library.plan(family, variant, slots)
                for region in plan.regions.values():
                    self.assertGreater(region.w, 0)
                    self.assertGreater(region.h, 0)
                    self.assertGreaterEqual(region.x, 0)
                    self.assertGreaterEqual(region.y, 0)
                    self.assertLessEqual(region.right, design.layout.width + 0.001)
                    self.assertLessEqual(region.bottom, design.layout.height + 0.001)

    def test_split_cover_supports_custom_slide_width(self):
        metrics = LayoutMetrics(width=10.0, height=7.5)
        plan = BlueprintLibrary(metrics).plan("cover", "split")
        self.assertLessEqual(plan.region("visual").right, metrics.width)

    def test_invalid_variants_raise(self):
        library = BlueprintLibrary(technical_design().layout)
        with self.assertRaises(LayoutError):
            library.plan("comparison", "typo")
        with self.assertRaises(LayoutError):
            library.plan("roadmap", "typo", slots=4)

    def test_same_semantics_produce_distinct_designs(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            temp = Path(temp_dir)
            technical = temp / "technical.pptx"
            editorial = temp / "editorial.pptx"
            first_report = build_sample(technical_design(), technical)
            second_report = build_sample(editorial_design(), editorial)

            self.assertFalse(first_report.errors)
            self.assertFalse(second_report.errors)
            self.assertNotEqual(technical.read_bytes(), editorial.read_bytes())

            with zipfile.ZipFile(technical) as archive:
                technical_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
            with zipfile.ZipFile(editorial) as archive:
                editorial_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")

            self.assertIn(technical_design().palette.primary, technical_xml)
            self.assertIn(editorial_design().palette.primary, editorial_xml)
            self.assertNotEqual(technical_xml, editorial_xml)

    def test_text_overflow_raises_before_save(self):
        compiler = SlideCompiler(technical_design())
        canvas = compiler.begin_slide(
            SlideFrame(title="Overflow check", number=1),
            family="custom",
        )
        with self.assertRaises(TextOverflowError):
            canvas.text(
                Region(1.0, 2.0, 1.0, 0.18),
                "This body text cannot fit in a tiny region at sixteen points.",
                role="body",
            )

    def test_no_wrap_text_checks_horizontal_overflow(self):
        compiler = SlideCompiler(technical_design())
        canvas = compiler.begin_slide(
            SlideFrame(title="No-wrap check", number=1),
            family="custom",
        )
        with self.assertRaises(TextOverflowError):
            canvas.text(
                Region(1.0, 2.0, 0.55, 0.6),
                "this_identifier_is_far_too_long",
                role="code",
                word_wrap=False,
            )

    def test_mixed_size_rich_text_uses_each_span_width(self):
        compiler = SlideCompiler(technical_design())
        canvas = compiler.begin_slide(
            SlideFrame(title="Rich text check", number=1),
            family="custom",
        )
        canvas.rich_text(
            Region(1.0, 2.0, 6.0, 1.2),
            [
                TextSpan("BIG ", role="metric", size_pt=48, color_role="primary"),
                TextSpan(
                    "small supporting explanation that should remain compact",
                    role="source",
                    size_pt=9,
                    color_role="muted_ink",
                ),
            ],
        )

    def test_code_block_rejects_lines_that_exceed_its_region(self):
        compiler = SlideCompiler(technical_design())
        canvas = compiler.begin_slide(
            SlideFrame(title="Code capacity", number=1),
            family="code",
            variant="full",
        )
        with self.assertRaises(TextOverflowError):
            canvas.code(["x = 1"] * 30)

    def test_full_code_layout_rejects_explanations(self):
        compiler = SlideCompiler(technical_design())
        canvas = compiler.begin_slide(
            SlideFrame(title="Full code", number=1),
            family="code",
            variant="full",
        )
        with self.assertRaises(LayoutError):
            canvas.code(
                ["x = 1"],
                [ContentItem("Explanation", "Must not be discarded.")],
            )

    def test_status_components_reserve_text_space(self):
        compiler = SlideCompiler(technical_design())
        process = compiler.begin_slide(
            SlideFrame(title="Process status", number=1),
            family="process",
            variant="rail",
            slots=2,
        )
        process.process(
            [
                ContentItem(
                    "Build",
                    "This detail must end before the status chip.",
                    "primary",
                    "GA",
                ),
                ContentItem("Operate", "Second step", "secondary"),
            ]
        )
        record = compiler.records[-1]
        detail = next(text for text in record.text if text.text.startswith("This detail"))
        status = next(text for text in record.text if text.text == "GA")
        self.assertLessEqual(detail.region.bottom, status.region.y)

        custom = compiler.begin_slide(
            SlideFrame(title="Node status", number=2),
            family="custom",
        )
        custom.node(
            Region(1.0, 2.0, 4.0, 1.2),
            ContentItem("Long node title", "Node detail", "accent", "PREVIEW"),
        )
        record = compiler.records[-1]
        title = next(text for text in record.text if text.text == "Long node title")
        status = next(text for text in record.text if text.text == "PREVIEW")
        self.assertLessEqual(title.region.right, status.region.x)

    def test_trailing_newline_is_counted_during_fit(self):
        compiler = SlideCompiler(technical_design())
        canvas = compiler.begin_slide(
            SlideFrame(title="Trailing newline", number=1),
            family="custom",
        )
        with self.assertRaises(TextOverflowError):
            canvas.text(
                Region(1.0, 2.0, 4.0, 0.45),
                "Visible line\n",
                role="body",
            )

    def test_default_text_color_adapts_to_dark_surface(self):
        compiler = SlideCompiler(technical_design())
        canvas = compiler.begin_slide(
            SlideFrame(title="Dark surface", number=1),
            family="custom",
        )
        region = Region(1.0, 2.0, 4.0, 1.0)
        canvas.box(region, fill="surface_alt", line="border_inverse")
        canvas.text(
            region.inset(0.2),
            "Readable default text",
            role="body",
            background="surface_alt",
        )
        report = compiler.validate(min_layout_families=1)
        self.assertFalse(
            any("Low contrast" in issue.message for issue in report.errors)
        )

    def test_invalid_region_and_arrow_head_are_rejected(self):
        compiler = SlideCompiler(technical_design())
        canvas = compiler.begin_slide(
            SlideFrame(title="Geometry check", number=1),
            family="custom",
        )
        with self.assertRaises(LayoutError):
            canvas.box(Region(1.0, 2.0, -0.5, 1.0))
        with self.assertRaises(LayoutError):
            canvas.arrow(
                (2.0, 2.0),
                (technical_design().layout.width, 2.0),
            )

    def test_image_cover_and_contain_modes(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            temp = Path(temp_dir)
            image_path = temp / "sample.png"
            Image.new("RGB", (400, 200), (180, 80, 40)).save(image_path)
            compiler = SlideCompiler(technical_design())
            canvas = compiler.begin_slide(
                SlideFrame(title="Image placement", number=1),
                family="custom",
            )
            canvas.image(Region(1.0, 2.0, 3.0, 2.0), image_path, mode="cover")
            canvas.image(Region(5.0, 2.0, 3.0, 2.0), image_path, mode="contain")
            report = compiler.save(
                temp / "images.pptx",
                expected_slides=1,
                min_layout_families=1,
            )
            self.assertFalse(report.errors)

    def test_more_than_two_sources_are_rejected(self):
        compiler = SlideCompiler(technical_design())
        sources = [
            Source("Org", f"Source {index}", "2026-07-15")
            for index in range(3)
        ]
        with self.assertRaises(LayoutError):
            compiler.begin_slide(
                SlideFrame(title="Sources", source=sources),
                family="statement",
            )
        self.assertEqual(compiler.slide_count, 0)

    def test_strict_contrast_failure_is_an_error(self):
        compiler = SlideCompiler(technical_design())
        canvas = compiler.begin_slide(
            SlideFrame(title="Contrast check", number=1),
            family="custom",
        )
        canvas.text(
            Region(1.0, 2.0, 4.0, 0.8),
            "Invisible body text",
            role="body",
            color="surface",
            background="surface",
        )
        report = compiler.validate(min_layout_families=1)
        self.assertTrue(
            any("Low contrast" in issue.message for issue in report.errors)
        )

    def test_three_repeated_layouts_fail_validation(self):
        compiler = SlideCompiler(technical_design())
        for number in range(1, 4):
            compiler.begin_slide(
                SlideFrame(title=f"Statement {number}", number=number),
                family="statement",
                variant="centered",
            )
        report = compiler.validate(min_layout_families=1)
        self.assertTrue(
            any("repeats three times" in issue.message for issue in report.errors)
        )

    def test_dark_source_footer_is_plain_text(self):
        compiler = SlideCompiler(technical_design())
        canvas = compiler.begin_slide(
            SlideFrame(
                title="Dark source",
                number=1,
                background_role="canvas_alt",
                inverse=True,
                source=Source(
                    "Example",
                    "Dark source",
                    "2026-07-15",
                    "https://example.com/dark",
                ),
            ),
            family="statement",
        )
        hyperlinks = []
        for shape in canvas.slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.startswith("Source:"):
                        hyperlinks.append(run.hyperlink.address)
        self.assertEqual(hyperlinks, [None])


if __name__ == "__main__":
    unittest.main()
