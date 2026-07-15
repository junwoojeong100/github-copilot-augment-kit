from __future__ import annotations

import math
import unicodedata
from collections.abc import Sequence
from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from .layouts import BlueprintLibrary
from .models import (
    ContentItem,
    CompilerError,
    DeckValidationError,
    DesignDNA,
    LayoutError,
    LayoutPlan,
    Region,
    SlideFrame,
    SlideRecord,
    Source,
    TextOverflowError,
    TextRecord,
    TextSpan,
    ValidationIssue,
    ValidationReport,
    contrast_ratio,
    hex_to_rgb,
    required_layout_families,
)


def _rgb(value: str) -> RGBColor:
    return RGBColor(*hex_to_rgb(value))


def _weighted_text_width(text: str, size_pt: float) -> float:
    units = 0.0
    for character in text:
        if character == "\t":
            units += 2.2
        elif character.isspace():
            units += 0.32
        elif unicodedata.east_asian_width(character) in {"W", "F", "A"}:
            units += 1.0
        elif character.isupper():
            units += 0.64
        elif character in "ilI.,:;!'|":
            units += 0.3
        else:
            units += 0.55
    return units * size_pt


def _estimated_lines(
    text: str,
    width_in: float,
    size_pt: float,
    word_wrap: bool,
) -> int:
    available_points = max(1.0, width_in * 72)
    paragraphs = text.split("\n")
    lines = 0
    for paragraph in paragraphs:
        if not paragraph:
            lines += 1
            continue
        if not word_wrap:
            lines += 1
            continue
        estimated = _weighted_text_width(paragraph, size_pt) / available_points
        lines += max(1, math.ceil(estimated * 1.08))
    return lines


class SlideCanvas:
    """A slide plus its semantic layout plan."""

    def __init__(
        self,
        compiler: "SlideCompiler",
        slide,
        plan: LayoutPlan,
        frame: SlideFrame,
    ):
        self.compiler = compiler
        self.slide = slide
        self.plan = plan
        self.frame = frame

    def region(self, name: str, inset: float = 0) -> Region:
        region = self.plan.region(name)
        return region.inset(inset) if inset else region

    def box(
        self,
        region: str | Region,
        *,
        fill: str = "surface",
        line: str | None = "border",
        corner: str | None = None,
        line_width_pt: float | None = None,
        depth: str | None = None,
    ):
        return self.compiler.box(
            self.slide,
            self._resolve(region),
            fill=fill,
            line=line,
            corner=corner,
            line_width_pt=line_width_pt,
            depth=depth,
        )

    def text(
        self,
        region: str | Region,
        text: str,
        *,
        role: str = "body",
        color: str | None = None,
        background: str | None = None,
        bold: bool = False,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
        margin: float = 0.03,
        size_pt: float | None = None,
        min_size_pt: float | None = None,
        font: str | None = None,
        line_spacing: float = 1.05,
        word_wrap: bool = True,
    ):
        return self.compiler.text(
            self.slide,
            self._resolve(region),
            text,
            role=role,
            color=color,
            background=background,
            bold=bold,
            align=align,
            valign=valign,
            margin=margin,
            size_pt=size_pt,
            min_size_pt=min_size_pt,
            font=font,
            line_spacing=line_spacing,
            word_wrap=word_wrap,
        )

    def rich_text(
        self,
        region: str | Region,
        spans: Sequence[TextSpan],
        *,
        background: str | None = None,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
        margin: float = 0.03,
    ):
        return self.compiler.rich_text(
            self.slide,
            self._resolve(region),
            spans,
            background=background,
            align=align,
            valign=valign,
            margin=margin,
        )

    def chip(
        self,
        region: str | Region,
        text: str,
        *,
        fill: str = "primary",
        color: str | None = None,
        size_pt: float | None = None,
    ):
        return self.compiler.chip(
            self.slide,
            self._resolve(region),
            text,
            fill=fill,
            color=color,
            size_pt=size_pt,
        )

    def node(
        self,
        region: str | Region,
        item: ContentItem,
        *,
        fill: str = "surface",
    ):
        panel_fill = "surface_alt" if self.frame.inverse and fill == "surface" else fill
        return self.compiler.node(
            self.slide,
            self._resolve(region),
            item,
            fill=panel_fill,
            inverse=self.frame.inverse,
        )

    def arrow(
        self,
        start: tuple[float, float],
        end: tuple[float, float],
        *,
        color: str = "primary",
        width_pt: float | None = None,
    ):
        return self.compiler.arrow(
            self.slide,
            start,
            end,
            color=color,
            width_pt=width_pt,
        )

    def image(
        self,
        region: str | Region,
        path: Path,
        *,
        mode: str = "cover",
        line: str | None = None,
    ):
        return self.compiler.image(
            self.slide,
            self._resolve(region),
            path,
            mode=mode,
            line=line,
        )

    def process(self, items: Sequence[ContentItem]) -> None:
        for index, item in enumerate(items, 1):
            region = self.region(f"step_{index}", inset=0.05)
            panel_fill = "surface_alt" if self.frame.inverse else "surface"
            self.compiler.box(
                self.slide,
                region,
                fill=panel_fill,
                line=item.color_role,
            )
            badge_region = Region(region.x + 0.12, region.y + 0.13, 0.38, 0.38)
            self.compiler.badge(
                self.slide,
                badge_region,
                str(index),
                fill=item.color_role,
            )
            self.compiler.text(
                self.slide,
                Region(region.x + 0.58, region.y + 0.12, region.w - 0.7, 0.33),
                item.title,
                role="secondary",
                color=item.color_role,
                background=panel_fill,
                bold=True,
                valign=MSO_ANCHOR.MIDDLE,
            )
            if item.detail:
                detail_top = region.y + 0.58
                detail_bottom = region.bottom - (0.52 if item.status else 0.12)
                self.compiler.text(
                    self.slide,
                    Region(
                        region.x + 0.16,
                        detail_top,
                        region.w - 0.32,
                        detail_bottom - detail_top,
                    ),
                    item.detail,
                    role="secondary",
                    color="muted_ink" if not self.frame.inverse else "muted_inverse",
                    background=panel_fill,
                    align=PP_ALIGN.CENTER,
                )
            if item.status:
                self.compiler.chip(
                    self.slide,
                    Region(region.x + 0.16, region.bottom - 0.42, region.w - 0.32, 0.28),
                    item.status,
                    fill=item.color_role,
                )

    def layers(self, items: Sequence[ContentItem]) -> None:
        for index, item in enumerate(items, 1):
            region = self.region(f"layer_{index}")
            self.compiler.box(
                self.slide,
                region,
                fill="surface_alt" if self.frame.inverse else "surface",
                line=item.color_role,
            )
            self.compiler.text(
                self.slide,
                Region(region.x + 0.18, region.y + 0.1, 1.65, region.h - 0.2),
                item.title,
                role="label",
                color=item.color_role,
                background="surface_alt" if self.frame.inverse else "surface",
                bold=True,
                size_pt=14,
                valign=MSO_ANCHOR.MIDDLE,
            )
            self.compiler.text(
                self.slide,
                Region(region.x + 1.92, region.y + 0.08, region.w - 2.1, region.h - 0.16),
                item.detail,
                role="secondary",
                color="ink_inverse" if self.frame.inverse else "ink",
                background="surface_alt" if self.frame.inverse else "surface",
                bold=True,
                valign=MSO_ANCHOR.MIDDLE,
            )

    def comparison(
        self,
        left_title: str,
        left_items: Sequence[str],
        right_title: str,
        right_items: Sequence[str],
        *,
        left_color: str = "secondary",
        right_color: str = "primary",
    ) -> None:
        for name, title, items, color in (
            ("left", left_title, left_items, left_color),
            ("right", right_title, right_items, right_color),
        ):
            region = self.region(name)
            panel_fill = "surface_alt" if self.frame.inverse else "surface"
            self.compiler.box(self.slide, region, fill=panel_fill, line=color)
            self.compiler.chip(
                self.slide,
                Region(region.x + 0.2, region.y + 0.18, min(2.2, region.w - 0.4), 0.34),
                title,
                fill=color,
            )
            self.compiler.bullets(
                self.slide,
                Region(region.x + 0.26, region.y + 0.78, region.w - 0.52, region.h - 1.0),
                items,
                bullet_color=color,
                text_color="ink_inverse" if self.frame.inverse else "ink",
                background=panel_fill,
                role="secondary",
            )

    def matrix(
        self,
        headers: Sequence[str],
        rows: Sequence[Sequence[str]],
        *,
        column_widths: Sequence[float] | None = None,
        header_fill: str = "ink",
    ) -> None:
        if len(headers) > 5 or len(rows) + 1 > 6:
            raise LayoutError("Matrix blueprints are limited to 6 rows x 5 columns.")
        if any(len(row) != len(headers) for row in rows):
            raise LayoutError("Every matrix row must match the header width.")
        region = self.region("table")
        if column_widths is None:
            widths = [region.w / len(headers)] * len(headers)
        else:
            if len(column_widths) != len(headers):
                raise LayoutError("column_widths must match the header count.")
            total = sum(column_widths)
            widths = [region.w * width / total for width in column_widths]

        row_height = region.h / (len(rows) + 1)
        x = region.x
        for header, width in zip(headers, widths):
            cell = Region(x, region.y, width, row_height)
            self.compiler.box(self.slide, cell, fill=header_fill, line="border")
            self.compiler.text(
                self.slide,
                cell.inset(0.08),
                header,
                role="label",
                color="ink_inverse",
                background=header_fill,
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
            )
            x += width

        for row_index, row in enumerate(rows, 1):
            x = region.x
            if self.frame.inverse:
                fill = "surface_alt" if row_index % 2 else "canvas_alt"
                text_color = "ink_inverse"
            else:
                fill = "surface" if row_index % 2 else "canvas"
                text_color = "ink"
            for value, width in zip(row, widths):
                cell = Region(x, region.y + row_index * row_height, width, row_height)
                self.compiler.box(self.slide, cell, fill=fill, line="border")
                self.compiler.text(
                    self.slide,
                    cell.inset(0.1),
                    value,
                    role="secondary",
                    color=text_color,
                    background=fill,
                    bold=False,
                    valign=MSO_ANCHOR.MIDDLE,
                )
                x += width

    def code(
        self,
        lines: Sequence[str],
        explanations: Sequence[ContentItem] = (),
    ) -> None:
        code_region = self.region("code")
        if explanations and "explain" not in self.plan.regions:
            raise LayoutError(
                "Code explanations require the split variant with an explain region."
            )
        self.compiler.box(
            self.slide,
            code_region,
            fill="surface_alt",
            line="border_inverse",
            corner="rounded",
        )
        if not lines:
            raise LayoutError("Code blocks require at least one line.")
        line_height = (code_region.h - 0.4) / len(lines)
        if line_height <= 0:
            raise TextOverflowError("Code lines exceed the available code region.")
        for index, line in enumerate(lines, 1):
            y = code_region.y + 0.2 + (index - 1) * line_height
            self.compiler.text(
                self.slide,
                Region(code_region.x + 0.22, y, 0.42, line_height),
                f"{index:02d}",
                role="label",
                color="muted_inverse",
                background="surface_alt",
                align=PP_ALIGN.RIGHT,
                font=self.compiler.design.typography.mono_font,
                word_wrap=False,
            )
            self.compiler.text(
                self.slide,
                Region(code_region.x + 0.78, y, code_region.w - 1.0, line_height),
                line or " ",
                role="code",
                color="ink_inverse",
                background="surface_alt",
                font=self.compiler.design.typography.mono_font,
                word_wrap=False,
            )
        if explanations and "explain" in self.plan.regions:
            explain_region = self.region("explain")
            item_height = explain_region.h / len(explanations)
            for index, item in enumerate(explanations):
                region = Region(
                    explain_region.x,
                    explain_region.y + index * item_height,
                    explain_region.w,
                    item_height,
                )
                self.compiler.badge(
                    self.slide,
                    Region(region.x, region.y + 0.1, 0.38, 0.38),
                    str(index + 1),
                    fill=item.color_role,
                )
                self.compiler.text(
                    self.slide,
                    Region(region.x + 0.52, region.y + 0.04, region.w - 0.52, 0.32),
                    item.title,
                    role="secondary",
                    color=item.color_role,
                    background=self.frame.background_role,
                    bold=True,
                )
                self.compiler.text(
                    self.slide,
                    Region(region.x + 0.52, region.y + 0.4, region.w - 0.52, region.h - 0.46),
                    item.detail,
                    role="secondary",
                    color="muted_inverse" if self.frame.inverse else "muted_ink",
                    background=self.frame.background_role,
                )

    def roadmap(self, items: Sequence[ContentItem]) -> None:
        for index, item in enumerate(items, 1):
            region = self.region(f"stage_{index}", inset=0.04)
            panel_fill = "surface_alt" if self.frame.inverse else "surface"
            self.compiler.box(self.slide, region, fill=panel_fill, line=item.color_role)
            self.compiler.text(
                self.slide,
                Region(region.x + 0.14, region.y + 0.11, 0.46, 0.34),
                f"{index:02d}",
                role="label",
                color=item.color_role,
                background=panel_fill,
                bold=True,
                size_pt=14,
                font=self.compiler.design.typography.mono_font,
                valign=MSO_ANCHOR.MIDDLE,
                margin=0,
            )
            self.compiler.text(
                self.slide,
                Region(region.x + 0.14, region.y + 0.55, region.w - 0.28, 0.38),
                item.title,
                role="secondary",
                color=item.color_role,
                background=panel_fill,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            self.compiler.text(
                self.slide,
                Region(region.x + 0.16, region.y + 1.02, region.w - 0.32, region.h - 1.18),
                item.detail,
                role="secondary",
                color="muted_inverse" if self.frame.inverse else "muted_ink",
                background=panel_fill,
                align=PP_ALIGN.CENTER,
            )

    def _resolve(self, region: str | Region) -> Region:
        return self.region(region) if isinstance(region, str) else region


class SlideCompiler:
    """Theme-agnostic PowerPoint compiler driven by explicit Design DNA."""

    def __init__(
        self,
        design: DesignDNA,
        *,
        title: str = "",
        subject: str = "",
        author: str = "GitHub Copilot",
    ):
        self.design = design
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(design.layout.width)
        self.presentation.slide_height = Inches(design.layout.height)
        self.presentation.core_properties.title = title
        self.presentation.core_properties.subject = subject
        self.presentation.core_properties.author = author
        self.blueprints = BlueprintLibrary(
            design.layout,
            spacing=design.shapes.spacing,
        )
        self.records: list[SlideRecord] = []
        self._issues: list[ValidationIssue] = []
        self._record_by_slide_id: dict[int, SlideRecord] = {}

    @property
    def slide_count(self) -> int:
        return len(self.presentation.slides)

    def begin_slide(
        self,
        frame: SlideFrame,
        *,
        family: str,
        variant: str = "default",
        slots: int = 4,
        card_grid: bool | None = None,
        allow_repeat: bool = False,
    ) -> SlideCanvas:
        sources = frame.sources()
        if len(sources) > 2:
            raise LayoutError("A slide footer supports at most two sources.")
        plan = self.blueprints.plan(family, variant, slots)
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        record = SlideRecord(
            number=self.slide_count,
            title=frame.title,
            family=plan.family,
            variant=plan.variant,
            card_grid=plan.card_grid if card_grid is None else card_grid,
            allow_repeat=allow_repeat,
        )
        self.records.append(record)
        self._record_by_slide_id[id(slide)] = record
        self.box(
            slide,
            Region(0, 0, self.design.layout.width, self.design.layout.height),
            fill=frame.background_role,
            line=None,
            corner="square",
            depth="flat",
        )
        if frame.show_header:
            self._header(slide, frame)
        if frame.show_footer and sources:
            self._source(
                slide,
                sources,
                inverse=frame.inverse,
                background_role=frame.background_role,
            )
        return SlideCanvas(self, slide, plan, frame)

    def box(
        self,
        slide,
        region: Region,
        *,
        fill: str = "surface",
        line: str | None = "border",
        corner: str | None = None,
        line_width_pt: float | None = None,
        depth: str | None = None,
    ):
        self._assert_bounds(region)
        corner = corner or self.design.shapes.corner_style
        depth = depth or self.design.shapes.depth
        fill_hex = self.design.palette.resolve(fill)
        if depth == "shadow":
            shadow = Region(region.x + 0.05, region.y + 0.05, region.w, region.h)
            self._assert_bounds(shadow)
            self._shape(
                slide,
                shadow,
                fill=self.design.palette.border,
                line=None,
                corner=corner,
            )
        return self._shape(
            slide,
            region,
            fill=fill_hex,
            line=self.design.palette.resolve(line) if line else None,
            corner=corner,
            line_width_pt=line_width_pt,
        )

    def circle(
        self,
        slide,
        region: Region,
        *,
        fill: str,
        line: str | None = None,
        line_width_pt: float | None = None,
    ):
        self._assert_bounds(region)
        if abs(region.w - region.h) > 0.001:
            raise LayoutError("Circle regions must have equal width and height.")
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(region.x),
            Inches(region.y),
            Inches(region.w),
            Inches(region.h),
        )
        self._style_shape(shape, fill, line, line_width_pt)
        return shape

    def line(
        self,
        slide,
        start: tuple[float, float],
        end: tuple[float, float],
        *,
        color: str = "border",
        width_pt: float | None = None,
    ):
        self._assert_point(start)
        self._assert_point(end)
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(start[0]),
            Inches(start[1]),
            Inches(end[0]),
            Inches(end[1]),
        )
        connector.line.color.rgb = _rgb(self.design.palette.resolve(color))
        connector.line.width = Pt(width_pt or self.design.shapes.line_width_pt)
        return connector

    def arrow(
        self,
        slide,
        start: tuple[float, float],
        end: tuple[float, float],
        *,
        color: str = "primary",
        width_pt: float | None = None,
        head_in: float = 0.18,
    ):
        self._assert_point(start)
        self._assert_point(end)
        if head_in <= 0:
            raise LayoutError("Arrow head size must be positive.")
        dx = end[0] - start[0]
        dy = end[1] - start[1]
        length = math.hypot(dx, dy)
        if length <= head_in:
            raise LayoutError("Arrow is too short for its head.")
        ux, uy = dx / length, dy / length
        line_end = (end[0] - ux * head_in * 0.65, end[1] - uy * head_in * 0.65)
        head_region = Region(
            end[0] - head_in / 2,
            end[1] - head_in / 2,
            head_in,
            head_in,
        )
        self._assert_bounds(head_region)
        self.line(
            slide,
            start,
            line_end,
            color=color,
            width_pt=width_pt,
        )
        triangle = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(end[0] - head_in / 2),
            Inches(end[1] - head_in / 2),
            Inches(head_in),
            Inches(head_in),
        )
        triangle.rotation = math.degrees(math.atan2(dy, dx)) + 90
        triangle.fill.solid()
        triangle.fill.fore_color.rgb = _rgb(self.design.palette.resolve(color))
        triangle.line.fill.background()
        return triangle

    def image(
        self,
        slide,
        region: Region,
        path: Path,
        *,
        mode: str = "cover",
        line: str | None = None,
    ):
        self._assert_bounds(region)
        image_path = path.expanduser().resolve()
        if not image_path.is_file():
            raise FileNotFoundError(image_path)
        if mode not in {"cover", "contain"}:
            raise LayoutError("Image mode must be 'cover' or 'contain'.")

        with PILImage.open(image_path) as image:
            image_width, image_height = image.size
        image_ratio = image_width / image_height
        target_ratio = region.w / region.h

        if mode == "contain":
            if image_ratio >= target_ratio:
                width = region.w
                height = width / image_ratio
            else:
                height = region.h
                width = height * image_ratio
            left = region.x + (region.w - width) / 2
            top = region.y + (region.h - height) / 2
            picture = slide.shapes.add_picture(
                str(image_path),
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height),
            )
        else:
            picture = slide.shapes.add_picture(
                str(image_path),
                Inches(region.x),
                Inches(region.y),
                Inches(region.w),
                Inches(region.h),
            )
            if image_ratio > target_ratio:
                visible = target_ratio / image_ratio
                crop = (1 - visible) / 2
                picture.crop_left = crop
                picture.crop_right = crop
            elif image_ratio < target_ratio:
                visible = image_ratio / target_ratio
                crop = (1 - visible) / 2
                picture.crop_top = crop
                picture.crop_bottom = crop

        if line:
            picture.line.color.rgb = _rgb(self.design.palette.resolve(line))
            picture.line.width = Pt(self.design.shapes.border_width_pt)
        else:
            picture.line.fill.background()
        return picture

    def text(
        self,
        slide,
        region: Region,
        text: str,
        *,
        role: str = "body",
        color: str | None = None,
        background: str | None = None,
        bold: bool = False,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
        margin: float = 0.03,
        size_pt: float | None = None,
        min_size_pt: float | None = None,
        font: str | None = None,
        line_spacing: float = 1.05,
        word_wrap: bool = True,
    ):
        self._assert_bounds(region)
        if not text and role != "source":
            self._issues.append(
                ValidationIssue(
                    "warning",
                    "Empty text frame requested.",
                    self._record(slide).number,
                )
            )
        requested = size_pt or self.design.typography.size_for(role)
        minimum = min_size_pt or self.design.typography.minimum_for(role)
        fitted = self._fit_font_size(
            text,
            region,
            requested,
            minimum,
            margin,
            line_spacing,
            word_wrap,
        )
        font_name = font or self.design.typography.font_for(role)
        color_role = (
            color
            or (self._best_text_role(background) if background else "ink")
        )
        color_hex = self.design.palette.resolve(color_role)
        self._check_contrast(
            slide,
            role=role,
            size_pt=fitted,
            bold=bold,
            color=color_role,
            background=background,
        )

        shape = slide.shapes.add_textbox(
            Inches(region.x),
            Inches(region.y),
            Inches(region.w),
            Inches(region.h),
        )
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = word_wrap
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        text_frame.margin_left = Inches(margin)
        text_frame.margin_right = Inches(margin)
        text_frame.margin_top = Inches(margin)
        text_frame.margin_bottom = Inches(margin)
        text_frame.vertical_anchor = valign
        for index, line in enumerate(text.split("\n")):
            paragraph = (
                text_frame.paragraphs[0]
                if index == 0
                else text_frame.add_paragraph()
            )
            paragraph.text = line or " "
            paragraph.alignment = align
            paragraph.line_spacing = line_spacing
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(fitted)
                run.font.bold = bold
                run.font.color.rgb = _rgb(color_hex)
        self._record(slide).text.append(
            TextRecord(
                slide_number=self._record(slide).number,
                role=role,
                text=text,
                size_pt=fitted,
                region=region,
            )
        )
        return shape

    def rich_text(
        self,
        slide,
        region: Region,
        spans: Sequence[TextSpan],
        *,
        background: str | None = None,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
        margin: float = 0.03,
    ):
        self._assert_bounds(region)
        if not spans:
            raise LayoutError("rich_text requires at least one TextSpan.")
        base_sizes = [
            span.size_pt or self.design.typography.size_for(span.role)
            for span in spans
        ]
        minimum_scale = max(
            self.design.typography.minimum_for(span.role) / size
            for span, size in zip(spans, base_sizes)
        )
        scale = self._fit_rich_text_scale(
            spans,
            base_sizes,
            region,
            minimum_scale,
            margin,
        )
        shape = slide.shapes.add_textbox(
            Inches(region.x),
            Inches(region.y),
            Inches(region.w),
            Inches(region.h),
        )
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        text_frame.margin_left = Inches(margin)
        text_frame.margin_right = Inches(margin)
        text_frame.margin_top = Inches(margin)
        text_frame.margin_bottom = Inches(margin)
        text_frame.vertical_anchor = valign
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = align
        for span, base_size in zip(spans, base_sizes):
            rendered_size = round(base_size * scale, 1)
            run = paragraph.add_run()
            run.text = span.text
            run.font.name = span.font or self.design.typography.font_for(span.role)
            run.font.size = Pt(rendered_size)
            run.font.bold = span.bold
            run.font.color.rgb = _rgb(self.design.palette.resolve(span.color_role))
            self._check_contrast(
                slide,
                role=span.role,
                size_pt=rendered_size,
                bold=span.bold,
                color=span.color_role,
                background=background,
            )
            self._record(slide).text.append(
                TextRecord(
                    slide_number=self._record(slide).number,
                    role=span.role,
                    text=span.text,
                    size_pt=rendered_size,
                    region=region,
                )
            )
        return shape

    def bullets(
        self,
        slide,
        region: Region,
        items: Sequence[str],
        *,
        bullet_color: str = "primary",
        text_color: str = "ink",
        background: str | None = None,
        role: str = "body",
        gap: float = 0.08,
    ) -> None:
        if not items:
            return
        row_height = (region.h - gap * (len(items) - 1)) / len(items)
        for index, item in enumerate(items):
            y = region.y + index * (row_height + gap)
            dot = min(0.12, row_height * 0.25)
            self.circle(
                slide,
                Region(region.x, y + row_height * 0.4, dot, dot),
                fill=bullet_color,
            )
            self.text(
                slide,
                Region(region.x + dot + 0.14, y, region.w - dot - 0.14, row_height),
                item,
                role=role,
                color=text_color,
                background=background,
                valign=MSO_ANCHOR.MIDDLE,
            )

    def badge(
        self,
        slide,
        region: Region,
        text: str,
        *,
        fill: str = "primary",
        color: str | None = None,
    ):
        self.circle(slide, region, fill=fill)
        return self.text(
            slide,
            region,
            text,
            role="label",
            color=color or self._best_text_role(fill),
            background=fill,
            bold=True,
            size_pt=max(14, self.design.typography.label_pt),
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
            word_wrap=False,
        )

    def chip(
        self,
        slide,
        region: Region,
        text: str,
        *,
        fill: str = "primary",
        color: str | None = None,
        size_pt: float | None = None,
    ):
        self.box(
            slide,
            region,
            fill=fill,
            line=None,
            corner="rounded",
            depth="flat",
        )
        return self.text(
            slide,
            region,
            text,
            role="label",
            color=color or self._best_text_role(fill),
            background=fill,
            bold=True,
            size_pt=size_pt or max(14, self.design.typography.label_pt),
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
            word_wrap=False,
        )

    def node(
        self,
        slide,
        region: Region,
        item: ContentItem,
        *,
        fill: str = "surface",
        inverse: bool = False,
    ):
        shape = self.box(
            slide,
            region,
            fill=fill,
            line=item.color_role,
        )
        title_width = region.w - 0.3 - (1.18 if item.status else 0)
        self.text(
            slide,
            Region(region.x + 0.15, region.y + 0.1, title_width, 0.32),
            item.title,
            role="secondary",
            color=item.color_role,
            background=fill,
            bold=True,
        )
        if item.detail:
            self.text(
                slide,
                Region(region.x + 0.15, region.y + 0.48, region.w - 0.3, region.h - 0.6),
                item.detail,
                role="secondary",
                color="muted_inverse" if inverse else "muted_ink",
                background=fill,
            )
        if item.status:
            self.chip(
                slide,
                Region(region.right - 1.15, region.y + 0.12, 0.95, 0.28),
                item.status,
                fill=item.color_role,
            )
        return shape

    def validate(
        self,
        *,
        expected_slides: int | None = None,
        min_layout_families: int | None = None,
        max_card_grid_ratio: float = 0.35,
    ) -> ValidationReport:
        issues = list(self._issues)
        slide_count = self.slide_count
        if expected_slides is not None and slide_count != expected_slides:
            issues.append(
                ValidationIssue(
                    "error",
                    f"Expected {expected_slides} slides, found {slide_count}.",
                )
            )

        families = tuple(dict.fromkeys(record.family for record in self.records))
        minimum = (
            min_layout_families
            if min_layout_families is not None
            else required_layout_families(slide_count)
        )
        if len(families) < minimum:
            issues.append(
                ValidationIssue(
                    "error",
                    f"Layout diversity requires {minimum} families; found {len(families)}.",
                )
            )

        card_count = sum(record.card_grid for record in self.records)
        card_ratio = card_count / slide_count if slide_count else 0.0
        if card_ratio > max_card_grid_ratio:
            issues.append(
                ValidationIssue(
                    "error",
                    f"Card-grid ratio {card_ratio:.1%} exceeds {max_card_grid_ratio:.1%}.",
                )
            )

        for index in range(2, len(self.records)):
            group = self.records[index - 2 : index + 1]
            if (
                len({record.family for record in group}) == 1
                and not any(record.allow_repeat for record in group)
            ):
                issues.append(
                    ValidationIssue(
                        "error",
                        f"Layout family {group[0].family!r} repeats three times.",
                        group[-1].number,
                    )
                )

        for record in self.records:
            for text in record.text:
                if text.role == "body" and text.size_pt < 16:
                    issues.append(
                        ValidationIssue(
                            "error",
                            f"Body text is {text.size_pt:.1f}pt: {text.text[:60]!r}",
                            record.number,
                        )
                    )

        return ValidationReport(
            slide_count=slide_count,
            layout_families=families,
            card_grid_ratio=card_ratio,
            issues=issues,
        )

    def save(
        self,
        output: Path,
        *,
        expected_slides: int | None = None,
        min_layout_families: int | None = None,
        max_card_grid_ratio: float = 0.35,
        report_path: Path | None = None,
        strict: bool = True,
    ) -> ValidationReport:
        report = self.validate(
            expected_slides=expected_slides,
            min_layout_families=min_layout_families,
            max_card_grid_ratio=max_card_grid_ratio,
        )
        if report_path:
            report.write_json(report_path)
        if strict and report.errors:
            messages = "; ".join(
                f"slide {issue.slide}: {issue.message}"
                if issue.slide
                else issue.message
                for issue in report.errors
            )
            raise DeckValidationError(messages)
        output.parent.mkdir(parents=True, exist_ok=True)
        self.presentation.save(output)
        return report

    def _header(self, slide, frame: SlideFrame) -> None:
        layout = self.design.layout
        accent = frame.accent_role
        ink = "ink_inverse" if frame.inverse else "ink"
        muted = "muted_inverse" if frame.inverse else "muted_ink"
        self.box(
            slide,
            Region(layout.left - 0.08, layout.title_top, 0.08, 0.54),
            fill=accent,
            line=None,
            corner="square",
        )
        if frame.section:
            self.text(
                slide,
                Region(layout.left + 0.12, layout.title_top, 1.8, 0.28),
                frame.section.upper(),
                role="label",
                color=accent,
                background=frame.background_role,
                bold=True,
                size_pt=14,
                valign=MSO_ANCHOR.MIDDLE,
                word_wrap=False,
            )
        self.text(
            slide,
            Region(
                layout.left + 0.1,
                layout.title_top + 0.28,
                layout.safe_width - 0.55,
                layout.content_top
                - (layout.title_top + 0.28)
                - (0.3 if frame.subtitle else 0.06),
            ),
            frame.title,
            role="title",
            color=ink,
            background=frame.background_role,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0.01,
        )
        if frame.subtitle:
            self.text(
                slide,
                Region(
                    layout.left + 0.12,
                    layout.content_top - 0.26,
                    layout.safe_width - 1.0,
                    0.26,
                ),
                frame.subtitle,
                role="secondary",
                color=muted,
                background=frame.background_role,
            )
        if frame.number is not None:
            self.text(
                slide,
                Region(layout.width - layout.right - 0.42, layout.title_top, 0.42, 0.28),
                f"{frame.number:02d}",
                role="label",
                color=muted,
                background=frame.background_role,
                align=PP_ALIGN.RIGHT,
                font=self.design.typography.mono_font,
                word_wrap=False,
            )

    def _source(
        self,
        slide,
        sources: Sequence[Source],
        *,
        inverse: bool,
        background_role: str,
    ) -> None:
        layout = self.design.layout
        line_role = "border_inverse" if inverse else "border"
        text_role = "muted_inverse" if inverse else "muted_ink"
        self.line(
            slide,
            (layout.left, layout.footer_top),
            (layout.width - layout.right, layout.footer_top),
            color=line_role,
            width_pt=0.6,
        )
        label = " | ".join(source.label() for source in sources)
        shape = self.text(
            slide,
            Region(layout.left + 0.02, layout.footer_top + 0.08, layout.safe_width - 0.8, 0.24),
            f"Source: {label}",
            role="source",
            color=text_role,
            background=background_role,
            word_wrap=False,
            margin=0,
        )
        if not inverse and len(sources) == 1 and sources[0].url:
            paragraph = shape.text_frame.paragraphs[0]
            if paragraph.runs:
                paragraph.runs[0].hyperlink.address = sources[0].url
                paragraph.runs[0].font.color.rgb = _rgb(
                    self.design.palette.resolve(text_role)
                )

    def _shape(
        self,
        slide,
        region: Region,
        *,
        fill: str,
        line: str | None,
        corner: str,
        line_width_pt: float | None = None,
    ):
        kind = (
            MSO_SHAPE.RECTANGLE
            if corner == "square"
            else MSO_SHAPE.ROUNDED_RECTANGLE
        )
        shape = slide.shapes.add_shape(
            kind,
            Inches(region.x),
            Inches(region.y),
            Inches(region.w),
            Inches(region.h),
        )
        if corner != "square":
            try:
                shape.adjustments[0] = 0.06 if corner == "subtle" else 0.18
            except (IndexError, ValueError):
                pass
        self._style_shape(shape, fill, line, line_width_pt)
        return shape

    def _style_shape(
        self,
        shape,
        fill: str,
        line: str | None,
        line_width_pt: float | None,
    ) -> None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(self.design.palette.resolve(fill))
        if line:
            shape.line.color.rgb = _rgb(self.design.palette.resolve(line))
            shape.line.width = Pt(
                line_width_pt or self.design.shapes.border_width_pt
            )
        else:
            shape.line.fill.background()

    def _fit_font_size(
        self,
        text: str,
        region: Region,
        requested: float,
        minimum: float,
        margin: float,
        line_spacing: float,
        word_wrap: bool,
    ) -> float:
        available_width = region.w - margin * 2
        available_width_points = max(1.0, available_width * 72)
        available_height_points = max(1.0, (region.h - margin * 2) * 72)
        candidate = requested
        while candidate >= minimum - 0.001:
            lines = _estimated_lines(text, available_width, candidate, word_wrap)
            required_height = lines * candidate * line_spacing * 1.05
            horizontal_fits = word_wrap or all(
                _weighted_text_width(paragraph, candidate) <= available_width_points
                for paragraph in text.split("\n")
            )
            if horizontal_fits and required_height <= available_height_points:
                return round(candidate, 1)
            candidate -= 0.5
        raise TextOverflowError(
            f"Text cannot fit region {region} at >= {minimum}pt: {text[:100]!r}"
        )

    def _best_text_role(self, background: str) -> str:
        background_hex = self.design.palette.resolve(background)
        return max(
            ("ink", "ink_inverse"),
            key=lambda role: contrast_ratio(
                background_hex,
                self.design.palette.resolve(role),
            ),
        )

    def _check_contrast(
        self,
        slide,
        *,
        role: str,
        size_pt: float,
        bold: bool,
        color: str,
        background: str | None,
    ) -> None:
        if not background:
            return
        color_hex = self.design.palette.resolve(color)
        background_hex = self.design.palette.resolve(background)
        large_text = size_pt >= 18 or (bold and size_pt >= 14)
        minimum_contrast = (
            4.5
            if role in {"body", "code"} or not large_text
            else 3.0
        )
        ratio = contrast_ratio(background_hex, color_hex)
        if ratio < minimum_contrast:
            self._issues.append(
                ValidationIssue(
                    "error",
                    f"Low contrast {ratio:.2f}:1 for {role} text.",
                    self._record(slide).number,
                )
            )

    def _fit_rich_text_scale(
        self,
        spans: Sequence[TextSpan],
        base_sizes: Sequence[float],
        region: Region,
        minimum_scale: float,
        margin: float,
    ) -> float:
        available_width_points = max(1.0, (region.w - margin * 2) * 72)
        available_height_points = max(1.0, (region.h - margin * 2) * 72)
        scale = 1.0
        while scale >= minimum_scale - 0.001:
            line_widths = [0.0]
            line_sizes = [0.0]
            for span, base_size in zip(spans, base_sizes):
                rendered_size = base_size * scale
                chunks = span.text.split("\n")
                for index, chunk in enumerate(chunks):
                    line_widths[-1] += _weighted_text_width(
                        chunk,
                        rendered_size,
                    )
                    line_sizes[-1] = max(line_sizes[-1], rendered_size)
                    if index < len(chunks) - 1:
                        line_widths.append(0.0)
                        line_sizes.append(0.0)

            required_height = 0.0
            for width, size in zip(line_widths, line_sizes):
                line_size = size or min(base_sizes) * scale
                wrapped_lines = max(
                    1,
                    math.ceil(width / available_width_points * 1.08),
                )
                required_height += wrapped_lines * line_size * 1.05 * 1.05
            if required_height <= available_height_points:
                return round(scale, 3)
            scale -= 0.02
        raise TextOverflowError(
            f"Rich text cannot fit region {region} without crossing a minimum size."
        )

    def _assert_bounds(self, region: Region) -> None:
        layout = self.design.layout
        if (
            not all(
                math.isfinite(value)
                for value in (region.x, region.y, region.w, region.h)
            )
            or region.w <= 0
            or region.h <= 0
            or region.x < -0.001
            or region.y < -0.001
            or region.right > layout.width + 0.001
            or region.bottom > layout.height + 0.001
        ):
            raise LayoutError(f"Region exceeds slide bounds: {region}")

    def _assert_point(self, point: tuple[float, float]) -> None:
        layout = self.design.layout
        if not all(math.isfinite(value) for value in point) or not (
            0 <= point[0] <= layout.width and 0 <= point[1] <= layout.height
        ):
            raise LayoutError(f"Point exceeds slide bounds: {point}")

    def _record(self, slide) -> SlideRecord:
        try:
            return self._record_by_slide_id[id(slide)]
        except KeyError as error:
            raise CompilerError("Slide was not created by this compiler.") from error
